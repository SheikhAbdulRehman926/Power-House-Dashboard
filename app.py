# powerhouse_dashboard_app.py
# Streamlit dashboard for "GENERATION AND EXPENSE" workbook
# Light theme, responsive tabs, centered background logos, context-aware chart layout selectors,
# explicit Solar Savings calculations, and export-perfect light charts — now with 1-click PDF, DOCX, and PPT per tab.
# ───────────────────────────────────────────────────────────────────────────────
# CHANGE LOG (this version):
# • Adjustable logo size script IN CODE (no sidebar): edit LOGO_MAX_HEIGHT_PX / LOGO_MAX_WIDTH_PX.
# • Watermark size fixed to match screenshot (50% width, 24% gap): edit WM_WIDTH_VW / WM_GAP_VW if needed.
# • PETPAK Prod vs Energy: bar = PETPAK orange; line changed from yellow to GREEN only.
# • Gas “Total Gas & Rate”: distinct, contrasting colors.
# • Expenses: high-contrast LESCO vs GAS.
# • Energy Mix: units on y-axis, taller & denser grid; snapshot hides selector menus.
# • Per-tab export buttons now show SIDE-BY-SIDE (PDF • DOCX • PPT) at the bottom of every tab.
# • NEW: Per-tab “Download PowerPoint (PPT)” — uses your uploaded template; slides = title → one slide per chart with summary → conclusion.
# • FIX: Solar Savings PDF now includes the actual charts, not just header (section key match).
# • FIX: Gas Consumption renders only on its tab (no leakage).
# • NEW: Global color-pick options + series multiselect on all charts; exports reflect layout/series/colors exactly.
# • NEW: Per-tab “Download Word report (CEO/COO/Director)” — text-only narrative describing every chart on the tab.
# • NEW: Forecasting tab (before Comparison): choose metric, Monthly vs Yearly, Gaussian Naive Bayes-based forecast with value card, ranges,
#        and chart (with color/layout selectors) — matches app UI logic.

from __future__ import annotations
from typing import Dict, List, Optional, Tuple
from collections import defaultdict

# PIL Decompression Bomb Protection - Prevent large image errors
import PIL.Image
PIL.Image.MAX_IMAGE_PIXELS = 200000000  # Set limit to 200M pixels (higher than default 178M)

# Performance optimization - Caching functions will be defined after streamlit import

import base64
import io
import os

import numpy as np
import pandas as pd
import plotly.express as px
import plotly.graph_objects as go
import plotly.io as pio
import streamlit as st

# Performance optimization - Add caching and session state to prevent unnecessary reruns
@st.cache_data(ttl=300)  # Cache for 5 minutes
def load_data_cached(path):
    """Load data with caching to prevent unnecessary reloads"""
    return load_sheet(path)

@st.cache_data(ttl=600)  # Cache for 10 minutes
def get_available_years_cached(df):
    """Get available years with caching"""
    return sorted(df['Year'].unique(), reverse=True)

@st.cache_data(ttl=600)  # Cache for 10 minutes
def get_available_months_cached(df):
    """Get available months with caching"""
    return sorted(df['Month'].unique(), reverse=True)

# Initialize session state for tab performance
if 'current_tab' not in st.session_state:
    st.session_state.current_tab = 'Overview'

if 'tab_loaded' not in st.session_state:
    st.session_state.tab_loaded = set()

# Initialize session state for better performance
if 'data' not in st.session_state:
    st.session_state.data = None
if 'computed_results' not in st.session_state:
    st.session_state.computed_results = {}
if 'chart_cache' not in st.session_state:
    st.session_state.chart_cache = {}

# Performance optimization settings
st.set_page_config(
    page_title="PowerPlant Dashboard",
    layout="wide",
    initial_sidebar_state="collapsed"
)

# ───────────────────────────────────────────────────────────────────────────────
# VIEWER/ADMIN PUBLISH SYSTEM (password via Environment Variable)
# ───────────────────────────────────────────────────────────────────────────────
import glob, datetime, shutil

# Folder to store published file(s) and the pointer file
DATA_DIR = os.environ.get("POWERHOUSE_DATA_DIR", "data")
# The “pointer” filename that all viewers read
PUBLISH_BASENAME = os.environ.get("POWERHOUSE_PUBLISH_NAME", "latest.xlsx")
# Name of the Environment Variable that contains the admin password
ADMIN_PWD_ENVVAR = os.environ.get("POWERHOUSE_ADMIN_PASSWORD_NAME", "POWERHOUSE_ADMIN_PASSWORD")

def _ensure_data_dir() -> None:
    os.makedirs(DATA_DIR, exist_ok=True)

def _published_path() -> str:
    """Path of the always-shown file for viewers."""
    return os.path.join(DATA_DIR, PUBLISH_BASENAME)

def _get_latest_file_path() -> str | None:
    """
    If a published pointer exists, use it. Otherwise, pick the newest *.xlsx
    under DATA_DIR. Returns None if nothing exists.
    """
    _ensure_data_dir()
    pub = _published_path()
    if os.path.exists(pub):
        return pub
    xl = sorted(glob.glob(os.path.join(DATA_DIR, "*.xlsx")), key=os.path.getmtime)
    return xl[-1] if xl else None

def _publish_uploaded_xlsx(uploaded_file) -> str:
    """
    Save this upload with a timestamp AND copy to 'latest.xlsx' so all viewers
    immediately see it. Returns the published pointer path.
    """
    _ensure_data_dir()
    ts = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
    raw_path = os.path.join(DATA_DIR, f"upload_{ts}.xlsx")
    with open(raw_path, "wb") as f:
        f.write(uploaded_file.getbuffer())
    pub = _published_path()
    shutil.copyfile(raw_path, pub)
    return pub

# PDF libs
from reportlab.lib.pagesizes import A4, landscape
from reportlab.pdfgen import canvas
from reportlab.lib.utils import ImageReader

# PDF merger to keep the SAME page style as per-tab exports
try:
    from PyPDF2 import PdfReader, PdfWriter
    _PDF_MERGE_OK = True
except Exception:
    _PDF_MERGE_OK = False


# DOCX (Word) — installed name `python-docx`
try:
    from docx import Document
    from docx.shared import Pt
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    _DOCX_OK = True
except Exception:
    _DOCX_OK = False

# PPTX (PowerPoint) — installed name `python-pptx`
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt as PPTPt
    from pptx.enum.text import PP_ALIGN
    _PPTX_OK = True
except Exception:
    _PPTX_OK = False

# ───────────────────────────────────────────────────────────────────────────────
# PAGE CONFIG
# ───────────────────────────────────────────────────────────────────────────────
st.set_page_config(
    page_title="PowerPlant Dashboard • PETPAK & GPAK",
    page_icon="ChatGPT Image Sep 8, 2025, 11_11_58 AM.png",  # <-- point to your image file
    layout="wide",
    initial_sidebar_state="expanded",
)

# ───────────────────────────────────────────────────────────────────────────────
# LOGO/WATERMARK SIZING — EDIT HERE (NO SIDEBAR CONTROLS)
# ───────────────────────────────────────────────────────────────────────────────
LOGO_MAX_HEIGHT_PX = 200  # cap logo height
LOGO_MAX_WIDTH_PX  = 230  # cap logo width
HDR_HEIGHT_PX      = 120  # fixed header height
HDR_PAD_V_PX       = 12   # vertical padding inside header
LOGO_SCALE         = 1.5  # keep 1.0 unless you want zoom (risk of clipping)

# watermark settings
WM_WIDTH_VW        = 50    # width of each watermark (in viewport %)
WM_GAP_VW          = 35    # horizontal gap between watermarks (in viewport %)
EXPORT_CLEAN       = True  # hide selectors in exported snapshots

# Global Plotly export config
CFG_PLOTLY = {
    "displaylogo": False,
    "modeBarButtonsToAdd": ["toImage"],
    "toImageButtonOptions": {"format": "png", "scale": 3, "filename": "powerhouse_chart"},
}

# ───────────────────────────────────────────────────────────────────────────────
# SECTION CAPTURE (for PDF, DOCX, PPT)
# ───────────────────────────────────────────────────────────────────────────────
SECTION_FIGS: Dict[str, List[go.Figure]] = defaultdict(list)
SECTION_CARDS: Dict[str, List[Dict]] = defaultdict(list)  # Store metric cards data
CURRENT_SECTION: List[str] = [""]  # tiny stack to allow nested usage if ever needed

def begin_section(name: str):
    """Set the current section; figures rendered after this will be captured."""
    CURRENT_SECTION[0] = name
    SECTION_FIGS[name] = []  # reset capture list on every render of a tab (fresh run)
    SECTION_CARDS[name] = []  # reset cards list on every render of a tab (fresh run)

def capture_fig(fig: go.Figure):
    name = CURRENT_SECTION[0] or "Page"
    SECTION_FIGS[name].append(fig)

def capture_metric_cards(cards_data: List[Dict], section_title: str = ""):
    """Capture metric cards data for PDF export.
    
    Args:
        cards_data: List of dictionaries with keys: label, value, delta, delta_color
        section_title: Title for the card section (e.g., "Yearly Highlights", "Monthly Highlights")
    """
    name = CURRENT_SECTION[0] or "Page"
    SECTION_CARDS[name].append({"cards": cards_data, "title": section_title})

# Cache figure rendering to prevent unnecessary regenerations
@st.cache_data(ttl=300)
def render_fig_cached(fig_dict, export_clean=False):
    """Render a cached Plotly figure to prevent unnecessary regenerations"""
    fig = go.Figure(fig_dict)
    if export_clean and fig.layout.updatemenus:
        d = fig.to_dict()
        d["layout"]["updatemenus"] = []
        fig = go.Figure(d)
    return fig

def render_fig(fig: go.Figure, key: str = None):
    """Consistent chart rendering + capture for PDF/DOCX/PPT. Optionally remove selectors."""
    # Generate unique key if not provided
    if key is None:
        import hashlib
        fig_str = str(fig.to_dict())
        key = f"chart_{hashlib.md5(fig_str.encode()).hexdigest()[:8]}"
    
    # Use cached rendering for better performance
    try:
        fig_dict = fig.to_dict()
        cached_fig = render_fig_cached(fig_dict, EXPORT_CLEAN)
        st.plotly_chart(cached_fig, use_container_width=True, config=CFG_PLOTLY, key=key)
    except Exception:
        # Fallback to direct rendering if caching fails
        if EXPORT_CLEAN and fig.layout.updatemenus:
            d = fig.to_dict()
            d["layout"]["updatemenus"] = []
            fig = go.Figure(d)
        st.plotly_chart(fig, use_container_width=True, config=CFG_PLOTLY, key=key)
    capture_fig(fig)

# ───────────────────────────────────────────────────────────────────────────────
# BRAND COLORS
# ───────────────────────────────────────────────────────────────────────────────
ORANGE_DARK   = "#C84B1A"
ORANGE_MAIN   = "#FF6A2C"
ORANGE_LIGHT  = "#FFA366"

BLUE_DARK     = "#0D2147"
BLUE_MAIN     = "#1E3A8A"
BLUE_LIGHT    = "#8EA6E8"

GAS_CRIMSON   = "#E11D48"  # Expenses GAS
LESCO_TEAL    = "#0EA5A7"  # Expenses LESCO
SOLAR_GOLD    = "#F59E0B"
RENTAL_GREEN  = "#22C55E"  # PETPAK production line (green)
PURPLE        = "#8B5CF6"
TEAL          = "#14B8A6"
CYAN          = "#06B6D4"
SLATE         = "#334155"
EMERALD       = "#10B981"  # Gas dual chart bar
VIOLET        = "#7C3AED"  # Gas dual chart line

def brand_palette(context: str, n: int) -> List[str]:
    if context == "petpak_engine":
        base = [ORANGE_MAIN, ORANGE_DARK, ORANGE_LIGHT]
    elif context == "gpak_engine_1":
        base = [BLUE_MAIN, BLUE_DARK, BLUE_LIGHT]
    elif context == "gpak_engine_2":
        base = [BLUE_LIGHT, BLUE_MAIN]
    elif context == "solar_pg":
        base = [ORANGE_MAIN, BLUE_MAIN]
    elif context == "lesco":
        base = [SLATE]
    elif context == "expenses":
        base = [LESCO_TEAL, GAS_CRIMSON, PURPLE, CYAN]
    elif context == "avg_cost":
        base = ["#DC2626"]
    elif context == "mix":
        base = [CYAN, TEAL, SOLAR_GOLD]  # LESCO, Gas, Solar
    else:
        base = [PURPLE, TEAL, ORANGE_MAIN, BLUE_MAIN, CYAN, SLATE]
    if n > len(base):
        base = (base * ((n // len(base)) + 1))[:n]
    else:
        base = base[:n]
    return base

# ───────────────────────────────────────────────────────────────────────────────
# LIGHT THEME
# ───────────────────────────────────────────────────────────────────────────────
TEXT_PRIMARY = "#0f172a"
TEXT_MUTED   = "#1f2937"
TAB_TEXT     = "#0f172a"
TAB_TEXT_SEL = "#0f172a"
CARD_BG      = "rgba(0,0,0,0.035)"

def _theme_text_color() -> str:
    return TEXT_PRIMARY

# ───────────────────────────────────────────────────────────────────────────────
# LOGOS
# ───────────────────────────────────────────────────────────────────────────────
def img_to_data_uri(path_list: List[str]) -> Optional[str]:
    for p in path_list:
        if os.path.exists(p):
            with open(p, "rb") as f:
                b64 = base64.b64encode(f.read()).decode("utf-8")
            return f"data:image/png;base64,{b64}"
    return None

PETPAK_LOGO = img_to_data_uri(
    ["/mnt/data/PETPAK-01.png", "PETPAK-01.png", "./PETPAK-01.png"]
)
GPAK_LOGO = img_to_data_uri(
    ["/mnt/data/GPAK-01.png", "GPAK-01.png", "./GPAK-01.png"]
)

# ───────────────────────────────────────────────────────────────────────────────
# CSS — THEME, TABS, WATERMARKS, HEADER
# ───────────────────────────────────────────────────────────────────────────────
left_logo_url  = GPAK_LOGO or ""
right_logo_url = PETPAK_LOGO or ""
pos_left  = f"calc(50% - {WM_GAP_VW}%)"
pos_right = f"calc(50% + {WM_GAP_VW}%)"

st.markdown(
    f"""
    <style>
    :root {{
        --text-primary:{TEXT_PRIMARY}; --text-muted:{TEXT_MUTED};
        --tab-text:{TAB_TEXT}; --tab-text-selected:{TAB_TEXT_SEL};
        --card-bg:{CARD_BG}; --orange-main:{ORANGE_DARK}; --blue-main:{BLUE_MAIN};
        --chip-border:rgba(0,0,0,0.10); --grid-light:#e5e7eb; 
        --logo-h:{LOGO_MAX_HEIGHT_PX}px;
        --logo-w:{LOGO_MAX_WIDTH_PX}px;
        --logo-scale:{LOGO_SCALE};
        --hdr-h:{HDR_HEIGHT_PX}px;
        --hdr-pad-v:{HDR_PAD_V_PX}px; 
        --wm-w:{WM_WIDTH_VW}vw;
        --wm-gap-left:{pos_left};
        --wm-gap-right:{pos_right};
    }}

    html,body,[data-testid="stAppViewContainer"],.main,.block-container {{ background:#fff!important;color:var(--text-primary); }}
    [data-testid="stHeader"] {{ background:transparent!important; height:18px; padding-top:2px; }}
    [data-testid="stDecoration"], [data-testid="stStatusWidget"] {{ display:none!important; }}

    .block-container::before {{
        content:""; position:fixed; top:50%; left:50%; transform:translate(-50%,-50%);
        width:min(92vw,1400px); height:min(52vw,760px); pointer-events:none; opacity:.06;
        background-image:url('{left_logo_url}'), url('{right_logo_url}');
        background-repeat:no-repeat,no-repeat;
        background-position:var(--wm-gap-left) center, var(--wm-gap-right) center;
        background-size:var(--wm-w) auto, var(--wm-w) auto; z-index:0;
    }}
    @media (max-width:1100px) {{
      .block-container::before {{
        background-position:center calc(50% - 22%), center calc(50% + 22%);
        background-size:min(60vw,420px) auto, min(60vw,420px) auto;
        width:100vw; height:min(100vh,1000px);
      }}
    }}
    .block-container > :not(style) {{ position:relative; z-index:1; }}

    .hdr {{
        display:grid; grid-template-columns:1fr auto 1fr; align-items:center; gap:8px;
        height:var(--hdr-h); min-height:var(--hdr-h); max-height:var(--hdr-h);
        padding:var(--hdr-pad-v) 16px; box-sizing:border-box; overflow:hidden;
        border-radius:12px;
        background:linear-gradient(90deg,{ORANGE_DARK} 0%,{ORANGE_MAIN} 45%,{BLUE_MAIN} 55%,{BLUE_DARK} 100%);
        box-shadow:0 4px 16px rgba(0,0,0,.18); margin-bottom:12px;
    }}
    .hdr .slot-left,.hdr .slot-right {{ display:flex; align-items:center; }}
    .hdr .slot-center {{ text-align:center; }}
    .hdr h1 {{
       color:#fff; font-weight:900; letter-spacing:.5px; margin:0;
       font-size:clamp(1.6rem, calc(1.6rem + 1.4vw), 3.0rem);
       line-height:1.3;
       white-space:nowrap; overflow:hidden; text-overflow:ellipsis;
     }}
    .hdr img.brand {{
        display:block; height:auto;
        max-height:min(var(--logo-h), calc(var(--hdr-h) - 2*var(--hdr-pad-v)));
        max-width:var(--logo-w);
        object-fit:contain;
        transform:scale(var(--logo-scale));
        transform-origin:center center;
        filter:drop-shadow(0 1px 2px rgba(0,0,0,.35));
    }}
    @media (max-width:680px) {{
      .hdr img.brand {{
        max-height:min(calc(var(--logo-h)*.8), calc(var(--hdr-h) - 2*var(--hdr-pad-v)));
        max-width:calc(var(--logo-w)*.8);
      }}
      .hdr .slot-left img.brand, .hdr .slot-right img.brand {{ opacity:.9; }}
    }}

    section[data-testid="stSidebar"] {{ background:#fff!important; border-right:1px solid rgba(0,0,0,.08)!important; }}

    .stTabs [role="tablist"] {{ display:flex!important; flex-wrap:nowrap!important; overflow-x:auto!important; gap:8px!important; margin-bottom:8px; }}
    .stTabs [role="tablist"]::-webkit-scrollbar {{ height:6px; }}
    .stTabs [role="tablist"]::-webkit-scrollbar-thumb {{ background:rgba(0,0,0,.2); border-radius:8px; }}
    .stTabs [role="tab"] {{
        flex:0 0 auto!important; border:1px solid rgba(0,0,0,.12); background:#fff; color:var(--tab-text);
        padding:10px 16px; border-radius:10px; font-weight:900; font-size:clamp(.90rem,.82rem + .28vw,1.04rem);
        text-transform:uppercase; letter-spacing:.6px; transition:transform .08s ease, background .15s ease;
    }}
    .stTabs [role="tab"]:hover {{ background:rgba(0,0,0,.03); transform:translateY(-1px); }}
    .stTabs [role="tab"][aria-selected="true"] {{
        background:linear-gradient(90deg,{ORANGE_MAIN} 0%, {BLUE_MAIN} 100%); color:#fff!important; border-color:transparent;
        text-shadow:0 1px 2px rgba(0,0,0,.35);
    }}

    .title-row {{ margin:6px 0 10px 0; }}
    .title-chip {{
      display:inline-flex; align-items:center; gap:.6rem; padding:.55rem .9rem; border-radius:12px; font-weight:800;
      color:#0f172a; background:#fff; border:1.5px solid var(--chip-border); box-shadow:0 1px 6px rgba(0,0,0,.06);
    }}
    .title-chip .dot {{ width:.65rem; height:.65rem; border-radius:999px;
      background:linear-gradient(180deg,{ORANGE_MAIN},{BLUE_MAIN}); box-shadow:0 0 0 2px rgba(0,0,0,.06) inset; }}

    div[data-testid="stMetric"] {{ position:relative; background:#fff; border-radius:12px; padding:12px 12px 12px 16px;
        box-shadow:inset 0 0 0 1px rgba(0,0,0,.08); color:var(--text-primary); }}
    div[data-testid="stMetric"]::before {{ content:""; position:absolute; left:0; top:6px; bottom:6px; width:6px;
        border-radius:6px; background:linear-gradient(180deg,{ORANGE_MAIN} 0%, {BLUE_MAIN} 100%); }}
    .element-container {{ box-shadow:none!important; background:transparent!important; }}
    .stDataFrame {{ background:#fff!important; border-radius:10px!important; box-shadow:inset 0 0 0 1px rgba(0,0,0,.08)!important; padding:6px!important; }}
    .js-plotly-plot,.plotly,.plot-container {{ color:var(--text-primary)!important; }}

   /* Export row: buttons styled like header */
    .export-row > div button, .export-row > div .stDownloadButton button {{
        width:100% !important;
        background:linear-gradient(90deg,{ORANGE_MAIN} 0%, {BLUE_MAIN} 100%) !important;
        color:#fff !important; font-weight:800 !important; border:none !important;
        border-radius:10px !important; box-shadow:0 2px 8px rgba(0,0,0,.18) !important;
        text-transform:uppercase; letter-spacing:.6px;
    }}
    .export-row > div button:hover, .export-row > div .stDownloadButton button:hover {{
        filter:brightness(1.03);
        transform:translateY(-1px);
    }}
    </style>
    """,
    unsafe_allow_html=True,
)

# ───────────────────────────────────────────────────────────────────────────────
# COMPACT CARDS TOGGLE
# ───────────────────────────────────────────────────────────────────────────────
compact_cards = st.toggle("Compact cards", value=True, key="ui_compact_cards")

# Add compact cards CSS if enabled
if compact_cards:
    st.markdown("""
    <style>
        /* Ultra-Tiny KPI Cards - Much Smaller */
        .metric-card {
            background: white;
            padding: 4px 6px;
            border-radius: 8px;
            border: 1px solid #e5e7eb;
            box-shadow: 0 1px 2px rgba(0,0,0,0.04);
            height: 40px;
            display: flex;
            flex-direction: column;
            justify-content: space-between;
            position: relative;
            overflow: hidden;
        }
        
        .metric-card::before {
            content: "";
            position: absolute;
            left: 0;
            top: 0;
            width: 6px;
            height: 100%;
            background: linear-gradient(to bottom, #ff6b35, #1e40af);
            border-radius: 14px 0 0 14px;
        }
        
        .metric-card .metric-title {
            font-size: 8px;
            font-weight: 500;
            color: #6b7280;
            margin-bottom: 8px;
            line-height: 1.2;
            margin-left: 8px;
        }
        
        .metric-card .metric-value {
            font-size: 12px;
            font-weight: 600;
            color: #1f2937;
            line-height: 1.2;
            margin-bottom: 12px;
            margin-left: 8px;
        }
        
        .metric-card .metric-delta {
            font-size: 7px;
            font-weight: 500;
            line-height: 1.0;
            margin-left: 8px;
        }
        
        .metric-card .metric-delta.positive {
            color: #059669;
        }
        
        .metric-card .metric-delta.negative {
            color: #dc2626;
        }
        
        /* Compact Controls */
        .compact-select {
            height: 36px;
            font-size: 12px;
        }
        
        .compact-select label {
            font-size: 12px;
            margin-bottom: 4px;
        }
        
        /* Cards Grid */
        .cards-grid {
            display: grid;
            grid-template-columns: repeat(4, 1fr);
            gap: 8px;
            margin-bottom: 20px;
        }
        
        @media (max-width: 768px) {
            .cards-grid {
                grid-template-columns: repeat(2, 1fr);
            }
        }
        
        @media (max-width: 480px) {
            .cards-grid {
                grid-template-columns: 1fr;
            }
        }
        
        /* Mega-Massive Charts - Even Bigger */
        .chart-container {
            background: white;
            padding: 2.5rem;
            border-radius: 1.25rem;
            border: 1px solid #e5e7eb;
            box-shadow: 0 8px 16px rgba(0,0,0,0.25);
            margin-bottom: 30px;
            min-height: 3000px;
        }
        
        .chart-container.stacked {
            min-height: 3600px;
        }
        
        .chart-container.comparison {
            min-height: 4200px;
        }
        
        .chart-container.single {
            min-height: 3300px;
        }
        
        /* Chart Typography */
        .chart-title {
            font-size: 18px;
            font-weight: 600;
            margin-bottom: 10px;
        }
        
        /* Swapped Layout - Charts Get Card Space, Cards Get Chart Space */
        .content-container {
            display: flex;
            flex-direction: column;
            gap: 30px;
        }
        
        .cards-section {
            flex: 0 0 auto;
            max-height: 5vh;
            overflow: hidden;
        }
        
        .charts-section {
            flex: 1;
            min-height: 3000px;
        }
        
        /* Proper spacing between cards and charts */
        .cards-to-charts-spacing {
            margin-bottom: 40px;
        }
        
        @media (max-width: 768px) {
            .charts-section {
                min-height: 3000px;
            }
            .cards-section {
                max-height: 8vh;
            }
        }
        
        @media (max-width: 480px) {
            .charts-section {
                min-height: 2500px;
            }
            .cards-section {
                max-height: 6vh;
            }
        }
    </style>
    """, unsafe_allow_html=True)

# ───────────────────────────────────────────────────────────────────────────────
# HEADER
# ───────────────────────────────────────────────────────────────────────────────
st.markdown(
    f"""
    <div class="hdr">
        <div class="slot-left">{f'<img class="brand" src="{GPAK_LOGO}" alt="GPAK">' if GPAK_LOGO else ""}</div>
        <div class="slot-center"><h1>Powerhouse Dashboard — PETPAK &amp; GPAK</h1></div>
        <div class="slot-right" style="justify-content:flex-end;">{f'<img class="brand" src="{PETPAK_LOGO}" alt="PETPAK">' if PETPAK_LOGO else ""}</div>
    </div>
    """,
    unsafe_allow_html=True,
)

# ───────────────────────────────────────────────────────────────────────────────
# FOOTER
# ───────────────────────────────────────────────────────────────────────────────
def render_footer_ui():
    st.markdown(
        f"""
        <style>
          .app-footer {{
            position: fixed; left: 0; right: 0; bottom: 0;
            height: 56px;
            background: linear-gradient(90deg, {ORANGE_DARK} 0%, {ORANGE_MAIN} 45%, {BLUE_MAIN} 55%, {BLUE_DARK} 100%);
            color: #fff; z-index: 1000; box-shadow: 0 -4px 16px rgba(0,0,0,.18);
            display: flex; align-items: center;
          }}
          .app-footer .row {{
            width: 100%; max-width: 1540px; margin: 0 auto; padding: 0 16px;
            display: flex; align-items: center; justify-content: space-between; gap: 12px;
            font-weight: 900; letter-spacing: .5px; white-space: nowrap;
          }}
          .block-container {{ padding-bottom: 80px; }}
        </style>
        <div class="app-footer">
          <div class="row">
            <div>PowerPlant Dashboard — PETPAK &amp; GPAK, version 1.0</div>
            <div>Powered by SARKS</div>
          </div>
        </div>
        """,
        unsafe_allow_html=True,
    )

# ───────────────────────────────────────────────────────────────────────────────
# HELPERS (color controls + export helpers + common charts)
# ───────────────────────────────────────────────────────────────────────────────
def section_title(text: str, *, level: int = 2):
    sizes = {1: "1.35rem", 2: "1.1rem", 3: "0.98rem"}
    size = sizes.get(level, sizes[2])
    st.markdown(
        f'''
        <div class="title-row">
          <div class="title-chip" style="font-size:{size}">
            <span class="dot"></span><span>{text}</span>
          </div>
        </div>
        ''',
        unsafe_allow_html=True,
    )

def _value_text(v) -> str:
    if isinstance(v, (int, float)) and not np.isnan(v):
        return f"{v:,.0f}" if abs(v) >= 1000 else f"{v:.2f}"
    return ""

def _add_point_labels(fig: go.Figure) -> None:
    for tr in fig.data:
        if getattr(tr, "type", "") == "scatter":
            tr.update(
                mode="lines+markers+text" if "lines" in (tr.mode or "") else "markers+text",
                text=[_value_text(y) for y in tr.y],
                textposition="top center",
                textfont=dict(color=TEXT_PRIMARY, size=11),
            )

def _add_bar_value_labels(fig: go.Figure, *, inside=False) -> None:
    color = "#FFFFFF" if inside else TEXT_PRIMARY
    for tr in fig.data:
        if getattr(tr, "type", "") == "bar":
            tr.update(
                texttemplate="<b>%{y:,.0f}</b>",
                textposition="inside" if inside else "outside",
                textfont=dict(size=11, color=color),
                insidetextanchor="middle",
                cliponaxis=False,
            )

def _apply_common_layout(fig: go.Figure, title: str) -> None:
    # Improved layout with better performance and visibility
    fig.update_layout(
        template="plotly_white",
        title=dict(
            text=title,
            x=0.5,
            xanchor='center',
            font=dict(size=16)
        ),
        hovermode="x unified",
        margin=dict(t=56, b=10, l=10, r=10, pad=4),
        legend=dict(
            orientation="h",
            y=-0.18,
            xanchor="center",
            x=0.5,
            itemclick="toggle",
            itemdoubleclick="toggleothers",
            bgcolor="rgba(255, 255, 255, 0.9)",
            bordercolor="rgba(0,0,0,0.1)",
            borderwidth=1
        ),
        yaxis=dict(
            autorange=True,
            tickformat="~s",
            showgrid=True,
            gridcolor="#e5e7eb",
            zeroline=True,
            zerolinecolor="#e5e7eb",
            zerolinewidth=1
        ),
        xaxis=dict(
            showgrid=True,
            gridcolor="#f1f5f9",
            zeroline=True,
            zerolinecolor="#f1f5f9",
            zerolinewidth=1
        ),
        plot_bgcolor="#ffffff",
        paper_bgcolor="#ffffff",
        font=dict(color=TEXT_PRIMARY, size=13),
        modebar=dict(bgcolor='rgba(255, 255, 255, 0.9)')
    )
    
    # Enable better performance for large datasets
    fig.update_traces(
        hovertemplate="<b>%{x}</b><br>" +
                     "%{y:,.2f}<br>" +
                     "<extra></extra>",
        hoverlabel=dict(
            bgcolor="white",
            font_size=13,
            font_family="Arial"
        )
    )

def _bar_like_from_wide(df: pd.DataFrame, x: str, y: str | List[str], title: str, *,
                        stacked=False, palette: Optional[List[str]] = None,
                        color_map: Optional[Dict[str, str]] = None) -> go.Figure:
    # Ensure y is a list for proper stacking
    y_cols = [y] if isinstance(y, str) else y
    
    fig = px.bar(df, x=x, y=y_cols, barmode=("stack" if stacked else "group"),
                 title=title, color_discrete_sequence=palette)
    
    if color_map:
        for tr in fig.data:
            if getattr(tr, "type", "") == "bar" and tr.name in color_map:
                tr.marker.color = color_map[tr.name]
    
    _apply_common_layout(fig, title)
    
    # Adjust label position based on stacking
    _add_bar_value_labels(fig, inside=(stacked and len(y_cols) > 1))
    
    # Ensure proper stacking order and layout
    if stacked:
        fig.update_layout(
            barmode='stack',
            bargap=0.15,
            bargroupgap=0.1
        )
    
    return fig

def _lollipop_from_series(df: pd.DataFrame, x: str, y: str, title: str, palette: Optional[List[str]] = None) -> go.Figure:
    fig = go.Figure()
    c = palette[0] if palette else None
    fig.add_bar(x=df[x], y=df[y], name=y, width=0.02, opacity=0.55, marker_color=c)
    _add_bar_value_labels(fig)
    fig.add_scatter(x=df[x], y=df[y], mode="markers+text", name="",
                    marker=dict(size=8, color=c),
                    text=[_value_text(v) for v in df[y]],
                    textposition="top center",
                    textfont=dict(color=TEXT_PRIMARY, size=11))
    _apply_common_layout(fig, title)
    return fig

def _line_from_wide(df: pd.DataFrame, x: str, ycols: List[str], title: str, palette: Optional[List[str]] = None) -> go.Figure:
    fig = go.Figure()
    palette = brand_palette("generic", len(ycols)) if palette is None else palette
    for i, c in enumerate(ycols):
        yv = pd.to_numeric(df[c], errors="coerce")
        fig.add_scatter(x=df[x], y=yv, name=c, mode="lines+markers+text",
                        line=dict(width=2.4, color=palette[i]),
                        marker=dict(size=6, color=palette[i]),
                        text=[_value_text(v) for v in yv],
                        textposition="top center",
                        textfont=dict(color=TEXT_PRIMARY, size=11),
                        hovertemplate="%{x|%b %Y}<br>%{y:,.4g}")
    _apply_common_layout(fig, title)
    return fig

def _area_from_wide(df: pd.DataFrame, x: str, ycols: List[str], title: str, palette: Optional[List[str]] = None) -> go.Figure:
    fig = go.Figure()
    palette = brand_palette("generic", len(ycols)) if palette is None else palette
    for i, c in enumerate(ycols):
        yv = pd.to_numeric(df[c], errors="coerce")
        fig.add_scatter(x=df[x], y=yv, name=c, fill="tozeroy", mode="lines+markers+text",
                        line=dict(width=2, color=palette[i]),
                        marker=dict(size=5, color=palette[i]),
                        text=[_value_text(v) for v in yv], textposition="top center",
                        textfont=dict(color=TEXT_PRIMARY, size=11),
                        hovertemplate="%{x|%b %Y}<br>%{y:,.4g}")
    _apply_common_layout(fig, title)
    return fig

def _overlay_totals_text(fig: go.Figure, x, totals: np.ndarray) -> int:
    idx = len(fig.data)
    fig.add_scatter(x=x, y=totals, mode="text", name="Total",
                    text=[f"<b>{_value_text(v)}</b>" for v in totals],
                    textfont=dict(color=TEXT_PRIMARY, size=11),
                    textposition="top center", hoverinfo="skip")
    return idx

def _add_visibility_dropdown_for_totals(fig: go.Figure, df: pd.DataFrame, ycols: List[str], totals_trace_idx: int):
    from itertools import combinations
    combos, names = [], []
    all_idx = list(range(len(ycols)))
    for r in range(1, len(ycols)+1):
        for cols in combinations(range(len(ycols)), r):
            combos.append(cols); names.append(" + ".join([ycols[i] for i in cols]))
    combos.append(tuple(all_idx)); names.append("All")

    buttons = []
    for cols, name in zip(combos, names):
        vis = [False] * (len(ycols) + 1)
        for i in cols: vis[i] = True
        sel_cols = [ycols[i] for i in cols]
        totals = df[sel_cols].sum(axis=1).values.astype(float)
        buttons.append(dict(label=name, method="update",
                            args=[{"visible": vis},
                                  {"data": [dict() for _ in range(len(ycols))] + [
                                      dict(y=totals, text=[f"<b>{_value_text(v)}</b>" for v in totals])
                                  ]}]))
    fig.update_layout(updatemenus=[dict(type="dropdown", direction="down",
                                        buttons=buttons, x=1.0, xanchor="right",
                                        y=1.15, yanchor="top", showactive=True,
                                        bgcolor="rgba(255,255,255,0.95)", pad={"r": 5, "t": 0})])

def safe_div(a, b):
    try:
        a = float(a); b = float(b)
    except (TypeError, ValueError):
        return np.nan
    if b == 0 or np.isnan(b):
        return np.nan
    return a / b

def _force_month_labels(fig: go.Figure):
    """
    Ensure Month axis shows as names like 'Jan 2025' instead of numbers/scientific notation.
    Used before exporting charts to PDF/PPT.
    """
    try:
        if fig.data and hasattr(fig.data[0], "x") and len(fig.data[0].x):
            xs = pd.to_datetime(pd.Series(fig.data[0].x), errors="coerce")
            if xs.notna().any():
                ticktext = xs.dt.strftime("%b %Y").tolist()
                for tr in fig.data:
                    if getattr(tr, "x", None) is not None and len(tr.x):
                        tr.x = ticktext
                fig.update_xaxes(
                    type="category",
                    tickmode="array",
                    tickvals=ticktext,
                    ticktext=ticktext,
                )
    except Exception:
        pass

def color_controls(labels: List[str], key_prefix: str, base_palette: List[str]) -> Tuple[List[str], Dict[str, str]]:
    """Global, consistent color control: checkbox + per-series color pickers."""
    use_custom = st.checkbox("Use custom colors", value=False, key=f"{key_prefix}_usecolors")
    color_map = {}
    palette = list(base_palette)
    if use_custom:
        cols = st.columns(min(3, len(labels)) if labels else 1)
        tmp_palette = []
        for i, lbl in enumerate(labels):
            with cols[i % len(cols)]:
                default = base_palette[i % len(base_palette)] if base_palette else "#4b5563"
                picked = st.color_picker(f"Color — {lbl}", value=default, key=f"{key_prefix}_color_{i}")
                color_map[lbl] = picked
                tmp_palette.append(picked)
        palette = tmp_palette
    return palette, color_map

# ───────────── DOCX narrative helpers (unchanged) ─────────────
def _series_stats(values: List[float]) -> Tuple[Optional[float], Optional[float], Optional[float], Optional[float]]:
    arr = pd.to_numeric(pd.Series(values), errors="coerce").dropna()
    if arr.empty: return (None, None, None, None)
    return float(arr.min()), float(arr.max()), float(arr.mean()), float(arr.sum())

def _fmt_val(v: Optional[float], suffix: str = "") -> str:
    if v is None or (isinstance(v, float) and np.isnan(v)):
        return "—"
    if abs(v) >= 1000 and suffix != "%":
        return f"{v:,.0f}{suffix}"
    if suffix == "%":
        return f"{v:.1f}%"
    return f"{v:.2f}{suffix}" if not float(v).is_integer() else f"{v:,.0f}{suffix}"

def _month_name(xval) -> str:
    try:
        dt = pd.to_datetime(xval)
        return dt.strftime("%b %Y")
    except Exception:
        return str(xval)

def _to_list_safe(v):
    if v is None:
        return []
    try:
        if isinstance(v, list):
            return v
        if isinstance(v, tuple):
            return list(v)
        try:
            import numpy as _np
            if isinstance(v, _np.ndarray):
                return v.tolist()
        except Exception:
            pass
        try:
            import pandas as _pd
            if isinstance(v, (_pd.Series, _pd.Index)):
                return v.tolist()
        except Exception:
            pass
        return list(v)
    except Exception:
        return [v] if v is not None else []

def _numeric_series(values):
    import pandas as _pd
    return _pd.to_numeric(_pd.Series(values), errors="coerce")

def _describe_trace(tr: go.BaseTraceType) -> str:
    ttype = getattr(tr, "type", "")
    name  = getattr(tr, "name", "") or ("Series" if ttype != "pie" else "Share")
    xs = _to_list_safe(getattr(tr, "x", None))
    ys = _to_list_safe(getattr(tr, "y", None))

    if ttype == "pie":
        labels = _to_list_safe(getattr(tr, "labels", None))
        values = _to_list_safe(getattr(tr, "values", None))
        if len(labels) and len(values):
            s_vals = _numeric_series(values)
            parts = []
            for lbl, val in zip(labels, list(s_vals)):
                parts.append(f"{lbl}: {_fmt_val(float(val))}" if pd.notna(val) else f"{lbl}: —")
            total = float(s_vals.sum(skipna=True)) if s_vals.notna().any() else None
            ttxt = _fmt_val(total) if total is not None else "—"
            return f"{name} shows composition — total {ttxt}; " + "; ".join(parts)
        return f"{name} shows composition data."

    if len(ys):
        s = _numeric_series(ys)
        if s.notna().any():
            mn = float(s.min(skipna=True)); mx = float(s.max(skipna=True)); avg = float(s.mean(skipna=True)); ssum = float(s.sum(skipna=True))
            try: idx_min = int(s.idxmin())
            except Exception: idx_min = None
            try: idx_max = int(s.idxmax())
            except Exception: idx_max = None

            def _m(i): return _month_name(xs[i]) if (i is not None and i < len(xs)) else "—"
            m_min = _m(idx_min); m_max = _m(idx_max)
            return (f"{name}: peak {_fmt_val(mx)} in {m_max}, lowest {_fmt_val(mn)} in {m_min}, "
                    f"average {_fmt_val(avg)}, total {_fmt_val(ssum)}.")
        else:
            return f"{name}: no numeric values available."
    return f"{name}: no numeric values available."

def _intro_for_section(title_text: str) -> str:
    return (f"Dear Sir,\n\n"
            f"This document summarizes insights from the **{title_text}** tab of the Powerhouse Dashboard "
            f"for PETPAK & GPAK. It provides a concise narrative of every chart and series displayed, "
            f"highlighting highs, lows, trends, and totals to support informed decision-making.\n")

def _outro_note() -> str:
    return ("\nIf you require deeper drill-downs (e.g., per asset, per shift, or rate sensitivity), kindly advise.\n\n"
            "Regards,\nTechnical Analytics Team")


def build_section_docx(section_name: str, title_text: str) -> bytes:
    if not _DOCX_OK:
        raise RuntimeError("python-docx not installed. Please run: pip install python-docx")

    figs = SECTION_FIGS.get(section_name, [])
    doc = Document()

    title = doc.add_paragraph()
    run = title.add_run(title_text)
    run.bold = True
    run.font.size = Pt(16)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER

    p = doc.add_paragraph(_intro_for_section(title_text))
    p.paragraph_format.space_after = Pt(8)

    for idx, fig in enumerate(figs, start=1):
        ftitle = (fig.layout.title.text if getattr(fig, "layout", None) and getattr(fig.layout, "title", None) else None) or f"Figure {idx}"
        h = doc.add_paragraph()
        r = h.add_run(f"{idx}. {ftitle}")
        r.bold = True
        r.font.size = Pt(12)

        traces = list(fig.data) if getattr(fig, "data", None) else []
        if not traces:
            doc.add_paragraph("• No chart data available for this figure.")
            continue
        for tr in traces:
            desc = _describe_trace(tr)
            doc.add_paragraph(f"• {desc}")
        doc.add_paragraph("• Where stacked bars are shown, combined monthly totals are labeled on top for quick comparison.")

    doc.add_paragraph(_outro_note())

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()

# ───────────── NEW: PPT (PowerPoint) builder & export row ─────────────
# ── PPT EXPORT (brand-styled, centered title, logos every slide, slide numbers) ─────────
from pptx import Presentation
from pptx.util import Inches, Pt as PPTPt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import io, os, numpy as np
import plotly.io as pio

PPT_SLIDE_W_IN = 13.333  # 1280px @ 96dpi
PPT_SLIDE_H_IN = 7.5
PPT_HEADER_H_IN = 0.9
PPT_MARGIN_X_IN = 0.5

def _rgb(hex_str: str):
    hs = hex_str.strip("#")
    return RGBColor(int(hs[0:2],16), int(hs[2:4],16), int(hs[4:6],16))

def _ppt_init() -> Presentation:
    prs = Presentation()
    prs.slide_width  = Inches(PPT_SLIDE_W_IN)
    prs.slide_height = Inches(PPT_SLIDE_H_IN)
    return prs

def _ppt_add_brand_header(slide, title_text: str = ""):
    # Left band (orange)
    left = slide.shapes.add_shape(1, Inches(0), Inches(0),
                                  Inches(PPT_SLIDE_W_IN/2), Inches(PPT_HEADER_H_IN))
    left.fill.solid(); left.fill.fore_color.rgb = _rgb(ORANGE_MAIN.strip())
    left.line.fill.background()
    # Right band (blue)
    right = slide.shapes.add_shape(1, Inches(PPT_SLIDE_W_IN/2), Inches(0),
                                   Inches(PPT_SLIDE_W_IN/2), Inches(PPT_HEADER_H_IN))
    right.fill.solid(); right.fill.fore_color.rgb = _rgb(BLUE_MAIN.strip())
    right.line.fill.background()
    # Optional title over the band
    if title_text:
        tb = slide.shapes.add_textbox(
            Inches(PPT_MARGIN_X_IN), Inches(0.12),
            Inches(PPT_SLIDE_W_IN - 2*PPT_MARGIN_X_IN), Inches(0.7)
        )
        tf = tb.text_frame; tf.clear()
        p = tf.paragraphs[0]; p.text = title_text
        p.font.size = PPTPt(26); p.font.bold = True
        p.font.color.rgb = RGBColor(255,255,255); p.alignment = PP_ALIGN.CENTER

def _ppt_add_logos(slide, petpak_logo_path: str = None, gpak_logo_path: str = None):
    # Always place logos on the header (every slide)
    y = 0.08; h = PPT_HEADER_H_IN - 0.16
    if PETPAK_LOGO and os.path.exists(PETPAK_LOGO):
        slide.shapes.add_picture(PETPAK_LOGO, Inches(0.2), Inches(y), height=Inches(h))
    elif petpak_logo_path and os.path.exists(petpak_logo_path):
        slide.shapes.add_picture(petpak_logo_path, Inches(0.2), Inches(y), height=Inches(h))
    if GPAK_LOGO and os.path.exists(GPAK_LOGO):
        slide.shapes.add_picture(GPAK_LOGO, Inches(PPT_SLIDE_W_IN - 0.2 - 1.8), Inches(y), height=Inches(h))
    elif gpak_logo_path and os.path.exists(gpak_logo_path):
        slide.shapes.add_picture(gpak_logo_path, Inches(PPT_SLIDE_W_IN - 0.2 - 1.8), Inches(y), height=Inches(h))

def _ppt_add_slide_number(slide, idx: int, total: int):
    # Bottom-right small, subtle
    w, h = 1.2, 0.3
    tb = slide.shapes.add_textbox(Inches(PPT_SLIDE_W_IN - w - 0.25),
                                  Inches(PPT_SLIDE_H_IN - h - 0.15),
                                  Inches(w), Inches(h))
    tf = tb.text_frame; tf.clear()
    p = tf.paragraphs[0]
    p.text = f"{idx} / {total}"
    p.font.size = PPTPt(11); p.font.color.rgb = RGBColor(100,116,139)  # slate-500
    p.alignment = PP_ALIGN.RIGHT

def _ppt_add_title_slide(prs: Presentation, title_text: str, slide_idx: int, total: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _ppt_add_brand_header(slide)   # header with band (no top title on band for cleaner hero)
    _ppt_add_logos(slide)
    # Centered hero title (both axes)
    tb = slide.shapes.add_textbox(
        Inches(1.2), Inches(PPT_SLIDE_H_IN/2 - 0.9),
        Inches(PPT_SLIDE_W_IN - 2.4), Inches(1.8)
    )
    tf = tb.text_frame; tf.clear()
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = PPTPt(40); p.font.bold = True
    p.font.color.rgb = RGBColor(30,41,59)  # slate-800
    p.alignment = PP_ALIGN.CENTER
    # Subtitle just below
    sb = slide.shapes.add_textbox(
        Inches(2.0), Inches(PPT_SLIDE_H_IN/2 + 0.5),
        Inches(PPT_SLIDE_W_IN - 4.0), Inches(0.8)
    )
    stf = sb.text_frame; stf.clear()
    sp = stf.paragraphs[0]
    sp.text = "Powerhouse Dashboard — PETPAK & GPAK"
    sp.font.size = PPTPt(20); sp.font.color.rgb = RGBColor(71,85,105)  # slate-600
    sp.alignment = PP_ALIGN.CENTER
    _ppt_add_slide_number(slide, slide_idx, total)

def _ppt_add_graph_slide(prs: Presentation, fig, caption: str, slide_idx: int, total: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _ppt_add_brand_header(slide, caption)
    _ppt_add_logos(slide)
    # Chart area
    top = PPT_HEADER_H_IN + 0.2
    width = PPT_SLIDE_W_IN - 2*PPT_MARGIN_X_IN
    _force_month_labels(fig)
    try:
        png = pio.to_image(fig, format="png", scale=2, width=1280, height=720)
    except Exception:
        png = pio.to_image(fig, format="png")
    slide.shapes.add_picture(io.BytesIO(png), Inches(PPT_MARGIN_X_IN), Inches(top), width=Inches(width))
    _ppt_add_slide_number(slide, slide_idx, total)

def _ppt_add_explain_slide(prs: Presentation, caption: str, bullets: list, slide_idx: int, total: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _ppt_add_brand_header(slide, f"{caption} — Explanation")
    _ppt_add_logos(slide)
    tb = slide.shapes.add_textbox(
        Inches(PPT_MARGIN_X_IN), Inches(PPT_HEADER_H_IN + 0.3),
        Inches(PPT_SLIDE_W_IN - 2*PPT_MARGIN_X_IN), Inches(PPT_SLIDE_H_IN - PPT_HEADER_H_IN - 0.6)
    )
    tf = tb.text_frame; tf.clear()
    if not bullets: bullets = ["No data available."]
    for i, t in enumerate(bullets[:14]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "• " + str(t)
        p.level = 0
        p.font.size = PPTPt(18); p.font.color.rgb = RGBColor(30,41,59)
    _ppt_add_slide_number(slide, slide_idx, total)

def _figure_summary_points(fig) -> list:
    """Professional, concise points per figure and trace."""
    pts = []
    title = getattr(getattr(fig, "layout", None), "title", None)
    if title and getattr(title, "text", None):
        pts.append(f"Overview: {title.text}")
    # For each trace, add last, delta, min/max, mean
    for tr in getattr(fig, "data", []):
        name = getattr(tr, "name", None) or "Series"
        try:
            y = np.asarray(tr.y, dtype=float)
            y = y[np.isfinite(y)]
        except Exception:
            y = np.array([])
        if y.size:
            last = y[-1]; mean = y.mean(); mn = y.min(); mx = y.max()
            delta = (y[-1] - y[-2]) if y.size > 1 else 0.0
            pts.append(
                f"{name}: latest {last:,.0f} (Δ {delta:+,.0f}); "
                f"avg {mean:,.0f}, range {mn:,.0f}–{mx:,.0f}."
            )
        else:
            pts.append(f"{name}: values shown; insufficient numeric history for stats.")
    # General recommendation seed
    if len(pts) < 2:
        pts.append("No significant variability detected; continue monitoring.")
    return pts

def _ppt_add_conclusion(prs: Presentation, points: list, slide_idx: int, total: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _ppt_add_brand_header(slide, "Conclusion & Recommendations")
    _ppt_add_logos(slide)
    tb = slide.shapes.add_textbox(
        Inches(PPT_MARGIN_X_IN), Inches(PPT_HEADER_H_IN + 0.3),
        Inches(PPT_SLIDE_W_IN - 2*PPT_MARGIN_X_IN), Inches(PPT_SLIDE_H_IN - PPT_HEADER_H_IN - 0.6)
    )
    tf = tb.text_frame; tf.clear()
    for i, t in enumerate(points or []):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "• " + str(t)
        p.level = 0
        p.font.size = PPTPt(18); p.font.color.rgb = RGBColor(30,41,59)
    _ppt_add_slide_number(slide, slide_idx, total)

def build_section_ppt_brand(section_name: str, title_text: str) -> bytes:
    if not _PPTX_OK:
        raise RuntimeError("python-pptx not installed. Please run: pip install python-pptx")

    figs = SECTION_FIGS.get(section_name, [])
    # total slides: 1 title + 2 per fig + 1 conclusion
    total = 1 + 2*len(figs) + 1
    idx = 1

    prs = _ppt_init()
    _ppt_add_title_slide(prs, title_text, idx, total); idx += 1

    for fig in figs:
        caption = (fig.layout.title.text
                   if getattr(fig, "layout", None) and getattr(fig.layout, "title", None)
                   else title_text)
        _ppt_add_graph_slide(prs, fig, caption, idx, total); idx += 1
        _ppt_add_explain_slide(prs, caption, _figure_summary_points(fig), idx, total); idx += 1

    _ppt_add_conclusion(prs, [
        "Sustain reliability by addressing sources with highest month-to-month variance.",
        "Prioritize solar utilization in months with favorable yield and tariff differentials.",
        "Investigate outliers in PKR/kWh and kWh/kg to identify maintenance or efficiency actions.",
        "Share this deck with operations; align next month’s targets to the median of the last 6 months."
    ], idx, total)

    buf = io.BytesIO(); prs.save(buf)
    return buf.getvalue()

# ── UTILITY FUNCTIONS FOR COST AND SAVINGS CALCULATIONS ──────────────────────────
def calc_costs_pkrs(solar_kwh: float, gas_kwh: float, lesco_kwh: float, rental_kwh: float) -> dict:
    """
    Calculate costs for different energy sources in PKR
    """
    # These rates are examples - adjust according to your actual rates
    RATES = {
        'solar': 15.5,  # PKR per kWh
        'gas': 18.2,    # PKR per kWh
        'lesco': 25.0,  # PKR per kWh
        'rental': 20.0  # PKR per kWh
    }
    
    costs = {
        'Solar Cost (PKR)': solar_kwh * RATES['solar'],
        'Gas Cost (PKR)': gas_kwh * RATES['gas'],
        'LESCO Cost (PKR)': lesco_kwh * RATES['lesco'],
        'Rental Cost (PKR)': rental_kwh * RATES['rental']
    }
    
    costs['Total Cost (PKR)'] = sum(costs.values())
    return costs

def compute_solar_savings(solar_kwh: float, baseline: str = "LESCO") -> float:
    """
    Calculate savings from solar usage compared to baseline source
    """
    # These rates are examples - adjust according to your actual rates
    RATES = {
        'LESCO': 25.0,  # PKR per kWh
        'GAS': 18.2,    # PKR per kWh
        'RENTAL': 20.0  # PKR per kWh
    }
    
    SOLAR_RATE = 15.5  # PKR per kWh
    
    if baseline not in RATES:
        raise ValueError(f"Invalid baseline source. Must be one of: {', '.join(RATES.keys())}")
    
    baseline_cost = solar_kwh * RATES[baseline]
    solar_cost = solar_kwh * SOLAR_RATE
    
    return baseline_cost - solar_cost

# ── REPORT EXPORT HELPERS (PDF/DOCX/PPT consolidated) ──────────────────────────
import io, os, numpy as np, time
import plotly.io as pio
import plotly.graph_objects as go

try:
    from reportlab.pdfgen import canvas as pdfcanvas
    from reportlab.lib.utils import ImageReader
    from reportlab.lib import colors
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
    _PDF_OK = True
except Exception:
    _PDF_OK = False

# Cache for data loading
@st.cache_data(ttl=3600)  # Cache for 1 hour
def load_data(file_path):
    return pd.read_excel(file_path)

# Cache for computations
@st.cache_data
def compute_statistics(df):
    return df.describe()

# Cache for visualization data
@st.cache_data
def prepare_chart_data(df, x_col, y_cols):
    return df.groupby(x_col)[y_cols].sum().reset_index()

try:
    from docx import Document
    from docx.shared import Inches as DocxInches, Pt as DocxPt
    _DOCX_OK = True
except Exception:
    _DOCX_OK = False

def _report_collect_figs():
    figs_report = SECTION_FIGS.get("Report", [])
    if figs_report:
        return [("Report", f) for f in figs_report]

    # Consolidated order (Overview → Comparison), EXCLUDING Forecasting
    section_order = [
        "Overview",
        "Energy Sources",
        "Solar Savings",
        "Expenses",
        "Production vs Consumption",
        "Gas Consumption",
        "Comparison",
    ]
    pairs = []
    for sec in section_order:
        for f in SECTION_FIGS.get(sec, []):
            pairs.append((sec, f))
    return pairs

def build_report_pdf_same_style(title_text: str) -> bytes:
    # Use a unique name to avoid scoping collisions
    consolidated_sections = [
        "Overview",
        "Energy Sources",
        "Solar Savings",
        "Expenses",
        "Production vs Consumption",
        "Gas Consumption",
        "Comparison",  # EXCLUDE "Forecasting"
    ]

    writer = PdfWriter()
    any_page = False

    for sec in section_order:
        # Skip empty sections gracefully
        if not SECTION_FIGS.get(sec):
            continue

        # Reuse your existing per-tab PDF builder → preserves exact style
        pdf_bytes = build_section_pdf(sec, f"{sec} — Powerhouse Dashboard", PETPAK_LOGO, GPAK_LOGO)

        # Append all pages from this section PDF
        reader = PdfReader(io.BytesIO(pdf_bytes))
        for page in reader.pages:
            writer.add_page(page)
            any_page = True

    if not any_page:
        raise RuntimeError("No charts found to export. Open tabs once, or ensure SECTION_FIGS is populated.")

    out = io.BytesIO()
    writer.write(out)
    return out.getvalue()

def build_report_docx(title_text: str) -> bytes:
    """
    Build a single DOCX that narrates ALL included sections (Overview → Comparison),
    explicitly excluding the Forecasting tab.
    """
    if not _DOCX_OK:
        raise RuntimeError("python-docx not installed. Please run: pip install python-docx")

    # Sections to include (no Forecasting)
    section_order = [
        "Overview",
        "Energy Sources",
        "Solar Savings",
        "Expenses",
        "Production vs Consumption",
        "Gas Consumption",
        "Comparison",
    ]

    doc = Document()

    # Title
    t = doc.add_paragraph()
    r = t.add_run(title_text)
    r.bold = True
    r.font.size = Pt(16)
    t.alignment = WD_ALIGN_PARAGRAPH.CENTER

    # Intro
    intro = _intro_for_section(title_text)
    p = doc.add_paragraph(intro)
    p.paragraph_format.space_after = Pt(8)

    # Sections
    for sec in section_order:
        figs = SECTION_FIGS.get(sec, [])
        if not figs:
            continue

        # Section heading
        sh = doc.add_paragraph()
        sr = sh.add_run(sec)
        sr.bold = True
        sr.font.size = Pt(14)

        # Figures
        for idx, fig in enumerate(figs, start=1):
            ftitle = (fig.layout.title.text
                      if getattr(fig, "layout", None) and getattr(fig.layout, "title", None)
                      else f"{sec} — Figure {idx}")
            # Figure header
            fh = doc.add_paragraph()
            fr = fh.add_run(f"{idx}. {ftitle}")
            fr.bold = True
            fr.font.size = Pt(12)

            traces = list(fig.data) if getattr(fig, "data", None) else []
            if not traces:
                doc.add_paragraph("• No chart data available for this figure.")
                continue

            # Bullets per trace
            for tr in traces:
                desc = _describe_trace(tr)
                doc.add_paragraph(f"• {desc}")

            # Generic note for stacked totals (helpful on bar charts)
            doc.add_paragraph("• For stacked bars, combined monthly totals are labeled on top where applicable.")

    # Outro
    doc.add_paragraph(_outro_note())

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def build_report_ppt_brand(title_text: str) -> bytes:
    """
    Build a branded PPT deck for ALL included sections (Overview → Comparison),
    explicitly excluding the Forecasting tab.
    """
    if not _PPTX_OK:
        raise RuntimeError("python-pptx not installed. Please run: pip install python-pptx")

    # Sections to include (no Forecasting)
    section_order = [
        "Overview",
        "Energy Sources",
        "Solar Savings",
        "Expenses",
        "Production vs Consumption",
        "Gas Consumption",
        "Comparison",
    ]

    # Flatten all figures in order
    figs = []
    for sec in section_order:
        sec_figs = SECTION_FIGS.get(sec, [])
        for f in sec_figs:
            figs.append((sec, f))

    if not figs:
        raise RuntimeError("No charts found to export. Open tabs once, or ensure SECTION_FIGS is populated.")

    # Slides: 1 title + (2 per figure) + 1 conclusion
    total = 1 + 2 * len(figs) + 1
    idx = 1

    prs = _ppt_init()

    # Title slide
    _ppt_add_title_slide(prs, title_text, idx, total)
    idx += 1

    # Per-figure slides
    for (sec, fig) in figs:
        caption = (fig.layout.title.text
                   if getattr(fig, "layout", None) and getattr(fig.layout, "title", None)
                   else f"{sec}")
        _ppt_add_graph_slide(prs, fig, caption, idx, total); idx += 1
        _ppt_add_explain_slide(prs, caption, _figure_summary_points(fig), idx, total); idx += 1

    # Conclusion
    _ppt_add_conclusion(prs, [
        "Sustain reliability by addressing sources with the highest variance.",
        "Prioritize solar utilization where yield and tariff differentials are favorable.",
        "Investigate cost outliers (₨/kWh) and kWh/kg to spot efficiency opportunities.",
        "Align next month’s targets to the median of the trailing six months."
    ], idx, total)

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()

def build_report_docx(title_text: str) -> bytes:
    """
    Build a single DOCX for ALL included sections (Overview → Comparison),
    explicitly EXCLUDING the Forecasting tab.
    """
    if not _DOCX_OK:
        raise RuntimeError("python-docx not installed. Please run: pip install python-docx")

    section_order = [
        "Overview",
        "Energy Sources",
        "Solar Savings",
        "Expenses",
        "Production vs Consumption",
        "Gas Consumption",
        "Comparison",
    ]

    doc = Document()

    # Title
    t = doc.add_paragraph()
    r = t.add_run(title_text)
    r.bold = True
    r.font.size = Pt(16)
    t.alignment = WD_ALIGN_PARAGRAPH.CENTER

    # Intro
    p = doc.add_paragraph(_intro_for_section(title_text))
    p.paragraph_format.space_after = Pt(8)

    # Sections & figures
    for sec in section_order:
        figs = SECTION_FIGS.get(sec, [])
        if not figs:
            continue

        sh = doc.add_paragraph()
        sr = sh.add_run(sec)
        sr.bold = True
        sr.font.size = Pt(14)

        for idx, fig in enumerate(figs, start=1):
            # Heading per figure
            ftitle = (fig.layout.title.text
                      if getattr(fig, "layout", None) and getattr(fig.layout, "title", None)
                      else f"{sec} — Figure {idx}")
            fh = doc.add_paragraph()
            fr = fh.add_run(f"{idx}. {ftitle}")
            fr.bold = True
            fr.font.size = Pt(12)

            # Narrative bullets per trace
            traces = list(getattr(fig, "data", []))
            if not traces:
                doc.add_paragraph("• No chart data available for this figure.")
                continue
            for tr in traces:
                desc = _describe_trace(tr)
                doc.add_paragraph(f"• {desc}")

            # Helpful note for stacked totals
            doc.add_paragraph("• For stacked bars, combined totals are labeled on top where applicable.")

    # Outro
    doc.add_paragraph(_outro_note())

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def build_report_ppt_brand(title_text: str) -> bytes:
    """
    Build a branded PPT for ALL included sections (Overview → Comparison),
    explicitly EXCLUDING the Forecasting tab.
    """
    if not _PPTX_OK:
        raise RuntimeError("python-pptx not installed. Please run: pip install python-pptx")

    section_order = [
        "Overview",
        "Energy Sources",
        "Solar Savings",
        "Expenses",
        "Production vs Consumption",
        "Gas Consumption",
        "Comparison",
    ]

    # Flatten figs in order
    figs = []
    for sec in section_order:
        for f in SECTION_FIGS.get(sec, []):
            figs.append((sec, f))

    if not figs:
        raise RuntimeError("No charts found to export. Open tabs once, or ensure SECTION_FIGS is populated.")

    # Slides: 1 title + (2 per figure) + 1 conclusion
    total = 1 + 2*len(figs) + 1
    idx = 1

    prs = _ppt_init()
    _ppt_add_title_slide(prs, title_text, idx, total); idx += 1

    for (sec, fig) in figs:
        caption = (fig.layout.title.text
                   if getattr(fig, "layout", None) and getattr(fig.layout, "title", None)
                   else sec)
        _ppt_add_graph_slide(prs, fig, caption, idx, total); idx += 1
        _ppt_add_explain_slide(prs, caption, _figure_summary_points(fig), idx, total); idx += 1

    _ppt_add_conclusion(prs, [
        "Sustain reliability by addressing sources with the highest variance.",
        "Prioritize solar utilization where yield and tariff differentials are favorable.",
        "Investigate cost outliers (₨/kWh) and kWh/kg to spot efficiency opportunities.",
        "Align next month’s targets to the median of the trailing six months."
    ], idx, total)

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()

# ───────────────────────────────────────────────────────────────────────────────

# Helper to draw a per-tab Download buttons in ONE ROW (PDF • DOCX • PPT)
def render_export_row(section_key: str, title_text: str, fname_slug: str):
    c1, c2, c3 = st.columns(3, gap="small")
    is_report = (section_key.strip().lower() == "report")

    with c1:
        try:
            pdf_bytes = build_report_pdf_same_style(title_text) if is_report \
                        else build_section_pdf(section_key, title_text, PETPAK_LOGO, GPAK_LOGO)
            st.download_button(
                label="Download PDF",
                data=pdf_bytes,
                file_name=f"{fname_slug}.pdf",
                mime="application/pdf",
                use_container_width=True,
                key=f"dl_{section_key}"
            )
        except RuntimeError as e:
            st.info(f"PDF export unavailable: {e}")

    with c2:
        try:
            docx_bytes = build_report_docx(title_text) if is_report \
                         else build_section_docx(section_key, title_text)
            st.download_button(
                label="Download DOCX",
                data=docx_bytes,
                file_name=f"{fname_slug}.docx",
                mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                use_container_width=True,
                key=f"docx_{section_key}",
            )
        except RuntimeError as e:
            st.info(f"DOCX export unavailable: {e}")

    with c3:
        try:
            ppt_bytes = build_report_ppt_brand(title_text) if is_report \
                        else build_section_ppt_brand(section_key, title_text)
            st.download_button(
                label="Download PPT",
                data=ppt_bytes,
                file_name=f"{fname_slug}.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                use_container_width=True,
                key=f"ppt_{section_key}",
            )
        except RuntimeError as e:
            st.info(f"PPT export unavailable: {e}")

    st.markdown('<div class="export-row"></div>', unsafe_allow_html=True)

# ───────────────────────────────────────────────────────────────────────────────
# PDF BUILDER (unchanged)
# ───────────────────────────────────────────────────────────────────────────────
def build_section_pdf(section_name: str, title_text: str,
                      petpak_logo_data_uri: Optional[str],
                      gpak_logo_data_uri: Optional[str]) -> bytes:
    figs = SECTION_FIGS.get(section_name, [])
    cards_data = SECTION_CARDS.get(section_name, [])
    buf = io.BytesIO()
    page_w, page_h = landscape(A4)
    c = canvas.Canvas(buf, pagesize=(page_w, page_h))

    # Margins & header area - optimized to reduce empty space
    M_L, M_R, M_T, M_B = 20, 20, 20, 20  # Reduced margins for more content space
    header_h = 60  # Reduced header height for more content space
    usable_w = page_w - M_L - M_R
    usable_h = page_h - M_T - M_B - header_h

    def data_uri_to_bytes(data_uri: Optional[str]) -> Optional[bytes]:
        if not data_uri:
            return None
        if data_uri.startswith("data:image"):
            try:
                b64 = data_uri.split(",", 1)[1]
                return base64.b64decode(b64)
            except Exception:
                return None
        try:
            if os.path.exists(data_uri):
                with open(data_uri, "rb") as f:
                    return f.read()
        except Exception:
            pass
        return None

    petpak_bytes = data_uri_to_bytes(petpak_logo_data_uri)
    gpak_bytes   = data_uri_to_bytes(gpak_logo_data_uri)

    c.setFillColorRGB(1, 1, 1)
    c.rect(0, page_h - header_h, page_w, header_h, fill=1, stroke=0)

    logo_max_h = 40
    logo_max_w = 120
    y_logo = page_h - 10 - logo_max_h
    if gpak_bytes:
        img = ImageReader(io.BytesIO(gpak_bytes))
        iw, ih = img.getSize()
        scale = min(logo_max_w / iw, logo_max_h / ih)
        c.drawImage(img, M_L, y_logo, iw * scale, ih * scale, mask='auto')
    if petpak_bytes:
        img = ImageReader(io.BytesIO(petpak_bytes))
        iw, ih = img.getSize()
        scale = min(logo_max_w / iw, logo_max_h / ih)
        c.drawImage(img, page_w - M_R - iw * scale, y_logo, iw * scale, ih * scale, mask='auto')

    c.setFont("Helvetica-Bold", 18)
    c.setFillColorRGB(0.09, 0.09, 0.12)
    c.drawCentredString(page_w / 2, page_h - 46, title_text)

    # Start content area
    current_y = page_h - header_h - 20
    
    # Render metric cards first
    if cards_data:
        for card_group_data in cards_data:
            if not card_group_data:
                continue
            
            # Extract cards and title
            if isinstance(card_group_data, dict) and 'cards' in card_group_data:
                cards = card_group_data['cards']
                card_section_title = card_group_data.get('title', '')
            else:
                # Backward compatibility for old format
                cards = card_group_data
                card_section_title = ''
            
            if not cards:
                continue
            
            # Draw section title if provided - dashboard style
            if card_section_title:
                c.setFont("Helvetica-Bold", 14)  # Dashboard section title size
                c.setFillColorRGB(0.1, 0.1, 0.1)
                c.drawString(M_L, current_y - 15, card_section_title)
                current_y -= 25  # Proper spacing for dashboard appearance
                
            # Calculate card dimensions - ultra-tiny design
            card_width = (usable_w - 30) / len(cards)  # 30 for spacing
            card_height = 40  # Ultra-tiny height
            
            # Draw cards - compact design matching screenshot
            for i, card in enumerate(cards):
                x = M_L + i * (card_width + 8)  # Reduced spacing (6-8px)
                y = current_y - card_height
                
                # Draw card background - compact style
                c.setFillColorRGB(1.0, 1.0, 1.0)  # Pure white background
                c.rect(x, y, card_width, card_height, fill=1, stroke=1)
                c.setFillColorRGB(0.9, 0.9, 0.9)  # Light gray border
                c.rect(x, y, card_width, card_height, fill=0, stroke=1)
                
                # Draw gradient accent bar on left side (6px width, orange→blue gradient)
                c.setFillColorRGB(1.0, 0.42, 0.21)  # Orange start
                c.rect(x, y, 6, card_height, fill=1, stroke=0)
                c.setFillColorRGB(0.12, 0.25, 0.69)  # Blue end
                c.rect(x, y + card_height/2, 6, card_height/2, fill=1, stroke=0)
                
                # Card content - compact typography
                c.setFillColorRGB(0.1, 0.1, 0.1)  # Dark text
                
                # Label - ultra-tiny style with increased spacing
                c.setFont("Helvetica", 8)  # Ultra-tiny label font
                label_text = card['label']
                # Replace currency symbols that might appear as black boxes
                label_text = label_text.replace('₨', 'Rs').replace('■', 'Rs')
                
                label_lines = label_text.split(' — ')
                if len(label_lines) > 1:
                    c.drawString(x + 8, y + card_height - 8, label_lines[0])  # Year/period
                    c.drawString(x + 8, y + card_height - 18, label_lines[1])  # Metric name (increased spacing)
                else:
                    c.drawString(x + 8, y + card_height - 12, label_text)  # Increased spacing
                
                # Value - ultra-tiny style with increased spacing
                c.setFont("Helvetica-Bold", 12)  # Ultra-tiny value font
                c.drawString(x + 8, y + card_height - 30, card['value'])  # Increased spacing from label
                
                # Delta (if exists) - ultra-tiny style
                if card.get('delta') and card['delta'] != '':
                    c.setFont("Helvetica", 7)  # Ultra-tiny delta font
                    delta_color = card.get('delta_color', 'normal')
                    if delta_color == 'inverse':
                        c.setFillColorRGB(0.1, 0.7, 0.1)  # Green for good changes
                    else:
                        c.setFillColorRGB(0.7, 0.1, 0.1)  # Red for bad changes
                    
                    # Add arrow based on delta
                    delta_text = card['delta']
                    if delta_text.startswith('+'):
                        delta_text = '▲ ' + delta_text
                    elif delta_text.startswith('-'):
                        delta_text = '▼ ' + delta_text
                    
                    c.drawString(x + 8, y + card_height - 38, delta_text)  # Adjusted for increased spacing
                    c.setFillColorRGB(0.1, 0.1, 0.1)  # Reset color
            
            current_y -= card_height + 10  # Minimal spacing - maximum space for charts
    
    # Render charts
    if figs:
        images: List[Tuple[int, int, bytes]] = []
        for fig in figs:
            _force_month_labels(fig)
            try:
                # Generate mega-massive charts - Even bigger PDF space utilization
                png = pio.to_image(fig, format="png", scale=7, width=6000, height=3600)
            except Exception:
                try:
                    png = pio.to_image(fig, format="png", scale=6, width=5500, height=3400)
                except Exception:
                    try:
                        png = pio.to_image(fig, format="png", scale=5, width=5000, height=3000)
                    except Exception:
                        try:
                            png = pio.to_image(fig, format="png", scale=4, width=4500, height=2800)
                        except Exception:
                            png = pio.to_image(fig, format="png", scale=3, width=4000, height=2500)
            # Add PIL decompression bomb protection
            try:
                img = ImageReader(io.BytesIO(png))
                iw, ih = img.getSize()
                images.append((iw, ih, png))
            except Exception as e:
                # If image is too large, try with smaller dimensions
                print(f"Warning: Chart too large, using fallback dimensions. Error: {e}")
                try:
                    # Fallback to smaller dimensions
                    png_fallback = pio.to_image(fig, format="png", scale=1, width=1200, height=700)
                    img = ImageReader(io.BytesIO(png_fallback))
                    iw, ih = img.getSize()
                    images.append((iw, ih, png_fallback))
                except Exception as fallback_error:
                    print(f"Error: Could not generate chart image. Error: {fallback_error}")
                    continue

        n = len(images)
        if n == 1:
            cols, rows = 1, 1
        elif n in (2, 3):
            cols, rows = 2, 2
        elif n == 4:
            cols, rows = 2, 2
        elif n in (5, 6):
            cols, rows = 3, 2
        elif n in (7, 8, 9):
            cols, rows = 3, 3
        else:
            cols, rows = 3, 3

        # Adjust available height for charts - fill MAXIMUM PDF space
        chart_area_h = current_y - M_B - 1  # Ultra-minimal margin to fill maximum space
        cell_w = (usable_w - (cols - 1) * 3) / cols  # Ultra-minimal spacing to maximize chart size
        cell_h = (chart_area_h - (rows - 1) * 3) / rows  # Ultra-minimal spacing to maximize chart size
        start_x = M_L
        start_y = current_y - 1  # Ultra-minimal margin to fill maximum space

        idx = 0
        for r in range(rows):
            y_top = start_y - r * (cell_h + 3)  # Ultra-minimal spacing to fill maximum space
            y_img = y_top - cell_h
            for ccol in range(cols):
                if idx >= n:
                    break
                iw, ih, png = images[idx]
                img = ImageReader(io.BytesIO(png))
                scale = min(cell_w / iw, cell_h / ih)
                draw_w, draw_h = iw * scale, ih * scale
                x_img = start_x + ccol * (cell_w + 3) + (cell_w - draw_w) / 2  # Ultra-minimal spacing to fill maximum space
                c.drawImage(img, x_img, y_img + (cell_h - draw_h) / 2, draw_w, draw_h, mask='auto')
                idx += 1

    c.showPage()
    c.save()
    return buf.getvalue()

# ───────────────────────────────────────────────────────────────────────────────
# DATA LOADING
# ───────────────────────────────────────────────────────────────────────────────
@st.cache_data(show_spinner=False)
def load_sheet(file_bytes_or_path, sheet_name: str = "GENERATION AND EXPENSE") -> pd.DataFrame:
    df = pd.read_excel(file_bytes_or_path, sheet_name=sheet_name, header=[0, 1], engine="openpyxl")

    new_cols = []
    for top, sub in df.columns:
        top_s, sub_s = str(top).strip(), str(sub).strip()
        if (top_s == "" or top_s.lower().startswith("unnamed")) and sub_s.lower() == "month":
            new_cols.append(("Month", "Month"))
        else:
            new_cols.append((top_s, sub_s))
    df.columns = pd.MultiIndex.from_tuples(new_cols)

    if ("Month", "Month") in df.columns:
        df[("Month", "Month")] = pd.to_datetime(df[("Month", "Month")], errors="coerce").dt.to_period("M").dt.to_timestamp()

    return df

def get_block(df: pd.DataFrame, block: str) -> Optional[pd.DataFrame]:
    cols = [c for c in df.columns if c[0] == block]
    if not cols:
        return None
    out = df[cols].copy()
    out.columns = [c[1] for c in cols]
    out.insert(0, "Month", df[("Month", "Month")])
    return out

def get_block_any(df: pd.DataFrame, names: List[str]) -> Optional[pd.DataFrame]:
    for name in names:
        blk = get_block(df, name)
        if blk is not None:
            return blk
    return None

def compute_rates(df_full: pd.DataFrame) -> Dict[str, float]:
    res = {}
    lesco_units = get_block(df_full, "LESCO GENERATION")
    if lesco_units is not None and "UNITS (KWH)" in lesco_units.columns:
        lesco_bill_col = ("EXPENSE", "LESCO BILL (Without GST)")
        if lesco_bill_col in df_full.columns:
            lesco_series_rate = df_full[lesco_bill_col] / lesco_units["UNITS (KWH)"].replace(0, np.nan)
            res["LESCO_PKRxKWh_avg"] = float(lesco_series_rate.mean(skipna=True))
    gas_units = get_block(df_full, "TOTAL GAS GENERATION")
    gas_bill_col = ("EXPENSE", "GAS BILL (Without GST)")
    if gas_units is not None and "UNITS (KWH)" in gas_units.columns and gas_bill_col in df_full.columns:
        gas_series_rate = df_full[gas_bill_col] / gas_units["UNITS (KWH)"].replace(0, np.nan)
        res["GAS_PKRxKWh_avg"] = float(gas_series_rate.mean(skipna=True))
    blended_col = ("EXPENSE", "PKR/KWh")
    if blended_col in df_full.columns:
        res["BLENDED_PKRxKWh_avg"] = float(df_full[blended_col].mean(skipna=True))
    return res

def compute_solar_savings_blended(df_full: pd.DataFrame) -> pd.DataFrame:
    """
    blended_rate_m = (LESCO_cost_m + GAS_cost_m) / (LESCO_kWh_m + GAS_kWh_m)
    savings_m      = Solar_kWh_m * blended_rate_m

    Returns columns your tab expects:
      Month, Solar_kWh, Saving_vs_GAS (now blended savings), GAS_rate_used (now blended rate)
    """
    if df_full is None or df_full.empty:
        return pd.DataFrame()

    # ---- helpers -------------------------------------------------------------
    def _pick_series(df: pd.DataFrame, *candidates):
        """
        Return the first non-empty numeric Series found among candidate columns.
        Each candidate can be a str (single-level column) or a tuple (MultiIndex).
        """
        for cand in candidates:
            try:
                if isinstance(df.columns, pd.MultiIndex) and isinstance(cand, tuple):
                    if cand in df.columns:
                        s = pd.to_numeric(df[cand], errors="coerce")
                    else:
                        continue
                else:
                    if cand in df.columns:
                        s = pd.to_numeric(df[cand], errors="coerce")
                    else:
                        continue
                # pick the first with at least one non-NaN value
                if s.notna().any():
                    return s
            except Exception:
                continue
        return None

    def _get_months(df: pd.DataFrame):
        if isinstance(df.columns, pd.MultiIndex) and ("Month", "Month") in df.columns:
            return pd.to_datetime(df[("Month", "Month")]).dt.to_period("M").dt.to_timestamp()
        if "Month" in df.columns:
            return pd.to_datetime(df["Month"]).dt.to_period("M").dt.to_timestamp()
        return None

    # ---- Solar kWh per month -------------------------------------------------
    solar = get_block(df_full, "SOLAR GENERATION")
    if solar is None:
        return pd.DataFrame()

    if "TOTAL" in solar.columns:
        solar_kwh = solar[["Month", "TOTAL"]].rename(columns={"TOTAL": "Solar_kWh"}).copy()
    else:
        cols = [c for c in ("PETPAK", "GPAK") if c in solar.columns]
        if not cols:
            return pd.DataFrame()
        solar_kwh = solar[["Month"] + cols].copy()
        solar_kwh["Solar_kWh"] = solar_kwh[cols].sum(axis=1)
        solar_kwh = solar_kwh[["Month", "Solar_kWh"]]

    # Normalize Month to month-start timestamps
    solar_kwh["Month"] = pd.to_datetime(solar_kwh["Month"]).dt.to_period("M").dt.to_timestamp()

    # ---- Units (kWh) per month ----------------------------------------------
    lesco_units = get_block(df_full, "LESCO GENERATION")
    gas_units   = get_block(df_full, "TOTAL GAS GENERATION")
    lesco_kwh = None
    gas_kwh   = None
    if lesco_units is not None and "UNITS (KWH)" in lesco_units.columns:
        lesco_kwh = lesco_units[["Month", "UNITS (KWH)"]].rename(columns={"UNITS (KWH)": "LESCO_kWh"}).copy()
        lesco_kwh["Month"] = pd.to_datetime(lesco_kwh["Month"]).dt.to_period("M").dt.to_timestamp()
    if gas_units is not None and "UNITS (KWH)" in gas_units.columns:
        gas_kwh = gas_units[["Month", "UNITS (KWH)"]].rename(columns={"UNITS (KWH)": "GAS_kWh"}).copy()
        gas_kwh["Month"] = pd.to_datetime(gas_kwh["Month"]).dt.to_period("M").dt.to_timestamp()

    # ---- Costs per month (from EXPENSE) --------------------------------------
    months_series = _get_months(df_full)
    # Fallback to solar months if Month column missing in df_full
    costs_months = months_series if months_series is not None else solar_kwh["Month"]

    # SAFE picking: no "or" on Series
    lesco_cost_s = _pick_series(
        df_full,
        ("EXPENSE", "LESCO BILL (Without GST)"),
        "LESCO BILL (Without GST)",
        ("EXPENSE", "LESCO BILL"),
        "LESCO BILL",
    )
    gas_cost_s = _pick_series(
        df_full,
        ("EXPENSE", "GAS BILL (Without GST)"),
        "GAS BILL (Without GST)",
        ("EXPENSE", "GAS BILL"),
        "GAS BILL",
    )

    costs = pd.DataFrame({"Month": costs_months})
    if lesco_cost_s is not None:
        costs["LESCO_cost"] = lesco_cost_s.values
    if gas_cost_s is not None:
        costs["GAS_cost"] = gas_cost_s.values

    # If EXPENSE rows are daily/transactional, aggregate to month
    costs["Month"] = pd.to_datetime(costs["Month"]).dt.to_period("M").dt.to_timestamp()
    costs = costs.groupby("Month", as_index=False).agg({"LESCO_cost": "sum", "GAS_cost": "sum"})

    # ---- Merge & compute blended rate ----------------------------------------
    out = solar_kwh.copy()
    if lesco_kwh is not None:
        out = out.merge(lesco_kwh, on="Month", how="left")
    if gas_kwh is not None:
        out = out.merge(gas_kwh, on="Month", how="left")
    out = out.merge(costs, on="Month", how="left")

    total_units  = (out.get("LESCO_kWh", 0).fillna(0) + out.get("GAS_kWh", 0).fillna(0)).replace(0, np.nan)
    total_costs  =  out.get("LESCO_cost", 0).fillna(0) + out.get("GAS_cost", 0).fillna(0)
    blended_rate = (total_costs / total_units).astype(float)         # ₨/kWh
    savings      = out["Solar_kWh"].fillna(0) * blended_rate         # ₨

    out = out[["Month", "Solar_kWh"]].copy()
    out["Saving_vs_GAS"] = savings          # kept name for UI compatibility
    out["GAS_rate_used"] = blended_rate     # kept name for UI compatibility
    return out

# ───────────────────────────────────────────────────────────────────────────────
# SIDEBAR — Upload only (no logo/watermark sliders)
# ───────────────────────────────────────────────────────────────────────────────
# ───────────────────────────────────────────────────────────────────────────────
# SIDEBAR — Mode toggle + Admin password unlock + Publish
# (Put BEFORE any code that needs df_full or creates tabs)
# ───────────────────────────────────────────────────────────────────────────────
with st.sidebar:
    st.title("Controls")
    st.subheader("Access")

    # Viewer/Admin toggle is always visible
    mode = st.radio("Select mode", ["Viewer", "Administrator"], index=0, key="app_mode")

    # Show which file is currently live for viewers
    current_path = _get_latest_file_path()
    if current_path and os.path.exists(current_path):
        ts = datetime.datetime.fromtimestamp(os.path.getmtime(current_path)).strftime("%Y-%m-%d %H:%M")
        st.caption(f"Currently showing: **{os.path.basename(current_path)}**  \nLast updated: {ts}")
    else:
        st.caption("No published data yet. Please publish from Administrator mode.")

    is_admin = False
    if mode == "Administrator":
        # Admin password expected from Environment Variable
        expected = os.environ.get(ADMIN_PWD_ENVVAR, "")
        pwd = st.text_input("Admin password", type="password", key="admin_pwd")

        if not expected:
            st.warning(
                f"Set the Environment Variable **{ADMIN_PWD_ENVVAR}** on the machine to enable Admin login."
            )
        elif pwd == expected and len(pwd) > 0:
            is_admin = True
            st.success("Admin unlocked.")

            # Only now show the uploader (browse option)
            up = st.file_uploader("Upload updated Excel (.xlsx)", type=["xlsx"], key="admin_upload")
            if up is not None and st.button("Publish for all viewers", use_container_width=True):
                pub = _publish_uploaded_xlsx(up)
                st.success(f"Published **{os.path.basename(pub)}**. Viewers will now see the updated dashboard.")
                # Remember this path for THIS session and refresh
                st.session_state["active_source_path"] = pub
                st.rerun()

        else:
            if len(pwd) > 0:
                st.error("Incorrect password.")

# Decide the active data source for THIS run
# Decide the active data source for THIS run (set earlier by viewer/admin gate)
active_path = st.session_state.get("active_source_path") or _get_latest_file_path()

# If nothing has been published yet, stop early
if not active_path or not os.path.exists(active_path):
    st.warning("No published Excel found. Switch to **Administrator** mode, enter the password, and upload the first file.")
    st.stop()

# ───────────────────────────────────────────────────────────────────────────────
# Load the Excel into df_full (your load_sheet knows the correct sheet)
# ───────────────────────────────────────────────────────────────────────────────
df_full = load_data_cached(active_path)
if df_full is None or (hasattr(df_full, "empty") and df_full.empty):
    st.error("Loaded file has no rows. Please check the Excel and re-upload in Administrator mode.")
    st.stop()

# ───────────────────────────────────────────────────────────────────────────────
# Robust Month normalization (works with duplicate/MultiIndex headers)
# ───────────────────────────────────────────────────────────────────────────────
def _pick_month_series(df: pd.DataFrame):
    """Return a 1-D Series for Month even if df['Month'] is a DataFrame
    (duplicate column names) or a MultiIndex first-level ('Month', <sub>)."""
    # 1) Plain single-level columns
    if "Month" in df.columns:
        s = df["Month"]
        if isinstance(s, pd.DataFrame):
            # pick first non-empty subcolumn
            for col in s.columns:
                cand = s[col]
                try:
                    if cand.notna().any():
                        return cand
                except Exception:
                    pass
            return s.iloc[:, 0]
        return s

    # 2) MultiIndex columns like ('Month', '…')
    for col in df.columns:
        if isinstance(col, tuple) and str(col[0]).strip().lower() == "month":
            return df[col]

    # 3) Case/space variations
    for col in df.columns:
        name = col[0] if isinstance(col, tuple) else col
        if str(name).strip().lower() == "month":
            return df[col]

    return None

_month_s = _pick_month_series(df_full)

if _month_s is not None:
    # Try your expected format first (e.g., "Jul-2025"), then fall back
    parsed = pd.to_datetime(_month_s, format="%b-%Y", errors="coerce")
    if parsed.isna().all():
        parsed = pd.to_datetime(_month_s, errors="coerce", dayfirst=True)

    # Write back to a single canonical column named 'Month'
    df_full["Month"] = parsed

    # If your code expects Month as index, uncomment this:
    # df_full = df_full.set_index("Month", drop=False)
else:
    st.warning("No 'Month' column found in the loaded sheet. Check header rows.")

# ───────────────────────────────────────────────────────────────────────────────
# EXTRACT BLOCKS
# ───────────────────────────────────────────────────────────────────────────────
blocks = {
    "PETPAK ENGINE #01": get_block(df_full, "PETPAK ENGINE #01"),
    "GPAK ENGINE #01": get_block(df_full, "GPAK ENGINE #01"),
    "GPAK ENGINE #2": get_block(df_full, "GPAK ENGINE #2"),
    "SOLAR GENERATION": get_block(df_full, "SOLAR GENERATION"),
    "LESCO GENERATION": get_block(df_full, "LESCO GENERATION"),
    "TOTAL GAS GENERATION": get_block(df_full, "TOTAL GAS GENERATION"),
    "TOTAL GENERATION": get_block(df_full, "TOTAL GENERATION"),
    "RENTAL ENGINE": get_block_any(df_full, [
        "RENTAL ENGINE", "RENTAL ENGINE #01", "RENTAL ENGINE #1",
        "RENTAL GENERATOR", "RENTAL GENSET", "PETPAK RENTAL ENGINE",
        "GPAK RENTAL ENGINE", "RENTAL ENGINE PETPAK", "RENTAL ENGINE GPAK"
    ]),
    "PETPAK PVC": get_block_any(df_full, [
        "PETPAK Goods Production vs Energy Consumption",
        "PETPAK Goods Production Vs. Energy Consumption",
        "PETPAK Goods Production Vs Energy Consumption",
    ]),
    "GPAK PVC": get_block_any(df_full, [
        "GPAK Goods Production vs Energy Consumption",
        "GPAK Goods Production Vs. Energy Consumption",
        "GPAK Goods Production Vs Energy Consumption",
    ]),
    "GAS CONSUMPTION": get_block(df_full, "GAS CONSUMPTION"),
}

def series(block_name: str, subcol: str = "UNITS (KWH)") -> Optional[pd.Series]:
    dfb = blocks.get(block_name)
    if dfb is None or subcol not in dfb.columns:
        return None
    return dfb.set_index("Month")[subcol].astype(float)

# Series
lesco_units   = series("LESCO GENERATION")
gas_total     = series("TOTAL GAS GENERATION")
gen_total     = series("TOTAL GENERATION")
petpak1_units = series("PETPAK ENGINE #01")
gpak1_units   = series("GPAK ENGINE #01")
gpak2_units   = series("GPAK ENGINE #2")
rental_units  = series("RENTAL ENGINE")

# Solar total
solar_total = None
if blocks["SOLAR GENERATION"] is not None:
    if "TOTAL" in blocks["SOLAR GENERATION"].columns:
        solar_total = blocks["SOLAR GENERATION"].set_index("Month")["TOTAL"].astype(float)
    else:
        cols_tmp = [c for c in ["PETPAK", "GPAK"] if c in blocks["SOLAR GENERATION"].columns]
        if cols_tmp:
            solar_total = blocks["SOLAR GENERATION"].set_index("Month")[cols_tmp].sum(axis=1).astype(float)

def pick_expense_col(df_full, sub_name_candidates: List[str]) -> Optional[pd.Series]:
    for cand in sub_name_candidates:
        key = ("EXPENSE", cand)
        if key in df_full.columns:
            return df_full[key]
    return None

expense_total = pick_expense_col(df_full, ["TOTAL", "TOTAL ", "Total"])
expense_lesco = pick_expense_col(df_full, ["LESCO BILL (Without GST)"])
expense_gas   = pick_expense_col(df_full, ["GAS BILL (Without GST)"])
pkr_per_kwh   = pick_expense_col(df_full, ["PKR/KWh"])

# Month axis & alignments
months = df_full[("Month", "Month")].dt.to_period("M").dt.to_timestamp()

def _align(s: Optional[pd.Series]) -> Optional[pd.Series]:
    if s is None:
        return None
    return s.reindex(months.values)

gen_aligned    = _align(gen_total)
exp_total_s    = pd.Series(expense_total.values, index=months) if expense_total is not None else None
pkwh_s         = pd.Series(pkr_per_kwh.values, index=months)   if pkr_per_kwh is not None else None
solar_aligned  = _align(solar_total)
lesco_aligned  = _align(lesco_units)
petpak1_al     = _align(petpak1_units)
gpak1_al       = _align(gpak1_units)
gpak2_al       = _align(gpak2_units)
rental_al      = _align(rental_units)

valid_mask = pd.Series(gen_aligned.notna().values, index=months.index) if gen_aligned is not None else pd.Series(False, index=months.index)
gen_mask = valid_mask
if gen_mask.any():
    last_month = months.loc[gen_mask].max()
else:
    last_month = months.max() if len(months) else None
selector_months = list(months.loc[gen_mask]) if gen_mask.any() else list(months.dropna())

# ───────────────────────────────────────────────────────────────────────────────
# TABS  (FORECASTING inserted before Comparison)
# ───────────────────────────────────────────────────────────────────────────────
tab_overview, tab_sources, tab_savings, tab_expense, tab_prodcons, tab_gas, tab_forecast, tab_compare, tab_data, tab_report = st.tabs(
    ["Overview", "Energy Sources", "Solar Savings", "Expenses", "Production vs Consumption", "Gas Consumption", "Forecasting", "Comparison", "Data", "Report"]
)

# ───────────────────────────────────────────────────────────────────────────────
# OVERVIEW
# ───────────────────────────────────────────────────────────────────────────────
with tab_overview:
    begin_section("Overview")
    section_title("Yearly Highlights", level=2)
    if selector_months:
        # Use cached year selection to prevent unnecessary reruns
        @st.cache_data(ttl=600)
        def get_years_cached(months):
            return sorted({d.year for d in months})
        
        years_available = get_years_cached(selector_months)
        default_year = last_month.year if last_month is not None else years_available[-1]
        sel_year = st.selectbox("Select year", options=years_available,
                                index=years_available.index(default_year),
                                key="ov_year_selector")
        prev_year = sel_year - 1 if (sel_year - 1) in years_available else None

        def _sum_masked(s: Optional[pd.Series], mask: pd.Series) -> float:
            if s is None: return np.nan
            vals = s.iloc[mask.values]
            return float(vals.sum(skipna=True)) if vals.notna().any() else np.nan

        def _avg_of_avg_cost(mask: pd.Series) -> float:
            if pkwh_s is not None:
                vals = pkwh_s.iloc[mask.values]
                vals = vals[vals.notna()]
                if len(vals): return float(vals.mean())
            tot_kwh = _sum_masked(gen_aligned, mask)
            tot_exp = _sum_masked(exp_total_s, mask)
            return tot_exp / tot_kwh if (not np.isnan(tot_kwh) and tot_kwh > 0 and not np.isnan(tot_exp)) else np.nan

        def _avg_of_solar_share(mask: pd.Series) -> float:
            if solar_aligned is None or gen_aligned is None: return np.nan
            s = (solar_aligned / gen_aligned) * 100.0
            vals = s.iloc[mask.values]
            vals = vals[vals.notna()]
            return float(vals.mean()) if len(vals) else np.nan

        y_mask  = (months.dt.year == sel_year) & gen_mask
        py_mask = (months.dt.year == prev_year) & gen_mask if prev_year else None

        y_kwh   = _sum_masked(gen_aligned, y_mask)
        y_exp   = _sum_masked(exp_total_s, y_mask)
        y_cost  = _avg_of_avg_cost(y_mask)
        y_share = _avg_of_solar_share(y_mask)

        if py_mask is not None and py_mask.any():
            p_kwh   = _sum_masked(gen_aligned, py_mask)
            p_exp   = _sum_masked(exp_total_s, py_mask)
            p_cost  = _avg_of_avg_cost(py_mask)
            p_share = _avg_of_solar_share(py_mask)
        else:
            p_kwh = p_exp = p_cost = p_share = np.nan

        f0   = lambda x: f"{x:,.0f}"
        f2   = lambda x: f"{x:.2f}"
        f1p  = lambda x: "—" if np.isnan(x) else f"{x:.1f}"
        fd0  = lambda d: f"{d:+,.0f}"
        fd2  = lambda d: f"{d:+,.2f}"
        fd1p = lambda d: "—" if np.isnan(d) else f"{d:+.1f}"

        c1, c2, c3, c4 = st.columns(4)
        yearly_cards = []
        
        def metric_delta_value(col, label, curr, prev, value_formatter, delta_formatter=None, inverse=False):
            def _fmt(v):
                return value_formatter(v) if v is not None and not (isinstance(v, float) and np.isnan(v)) else "—"
            if prev is None or (isinstance(prev, float) and np.isnan(prev)):
                col.metric(label=label, value=_fmt(curr), delta="")
                yearly_cards.append({"label": label, "value": _fmt(curr), "delta": "", "delta_color": "normal"})
                return
            if curr is None or (isinstance(curr, float) and np.isnan(curr)):
                col.metric(label=label, value="—", delta="")
                yearly_cards.append({"label": label, "value": "—", "delta": "", "delta_color": "normal"})
                return
            delta = curr - prev
            delta_txt = delta_formatter(delta) if delta_formatter else f"{delta:+,.0f}"
            delta_color = "inverse" if inverse else "normal"
            col.metric(label=label, value=_fmt(curr), delta=delta_txt, delta_color=delta_color)
            yearly_cards.append({"label": label, "value": _fmt(curr), "delta": delta_txt, "delta_color": delta_color})

        metric_delta_value(c1, f"{sel_year} — Total Generation (kWh)", y_kwh, p_kwh if prev_year else None, f0, fd0, inverse=False)
        metric_delta_value(c2, f"{sel_year} — Total Energy Expense (₨)", y_exp, p_exp if prev_year else None, f0, fd0, inverse=True)
        metric_delta_value(c3, f"{sel_year} — Average Cost (₨/kWh)", y_cost, p_cost if prev_year else None, f2, fd2, inverse=True)
        metric_delta_value(c4, f"{sel_year} — Solar Share (%)", y_share, p_share if prev_year else None, f1p, fd1p, inverse=False)
        
        # Capture yearly cards for PDF
        capture_metric_cards(yearly_cards, "Yearly Highlights")
    else:
        st.info("No dated rows found. Please ensure the 'Month' column is populated.")

    section_title("Monthly Highlights", level=2)
    if selector_months:
        default_idx = selector_months.index(last_month) if (last_month is not None and last_month in selector_months) else len(selector_months) - 1
        sel_month = st.selectbox("Select month", options=selector_months, index=default_idx,
                                 format_func=lambda d: d.strftime("%b %Y") if pd.notna(d) else "—",
                                 key="ov_month_selector")

        def _prev_valid(before: pd.Timestamp) -> Optional[pd.Timestamp]:
            prevs = [m for m in selector_months if m < before]
            return prevs[-1] if prevs else None
        prev_month = _prev_valid(sel_month) if sel_month is not None else None

        def _pt(s: Optional[pd.Series], when) -> Optional[float]:
            try:
                return float(s.loc[when]) if s is not None and when in s.index else np.nan
            except Exception:
                return np.nan

        kwh_curr = _pt(gen_aligned, sel_month)
        kwh_prev = _pt(gen_aligned, prev_month)
        exp_curr_val = _pt(exp_total_s, sel_month)
        exp_prev_val = _pt(exp_total_s, prev_month)
        cpk_curr_val = _pt(pkwh_s, sel_month)
        cpk_prev_val = _pt(pkwh_s, prev_month)
        solar_curr = _pt(solar_aligned, sel_month)
        solar_prev = _pt(solar_aligned, prev_month)
        mix_curr = safe_div(solar_curr, kwh_curr) * 100.0 if kwh_curr and not np.isnan(kwh_curr) else np.nan
        mix_prev = safe_div(solar_prev, kwh_prev) * 100.0 if kwh_prev and not np.isnan(kwh_prev) else np.nan

        f0, f2 = (lambda x: f"{x:,.0f}"), (lambda x: f"{x:.2f}")
        f1p = lambda x: "—" if np.isnan(x) else f"{x:.1f}"
        fd0, fd2, fd1p = (lambda d: f"{d:+,.0f}"), (lambda d: f"{d:+,.2f}"), (lambda d: "—" if np.isnan(d) else f"{d:+.1f}")

        cA, cB, cC, cD = st.columns(4)
        monthly_cards = []
        
        def metric_delta_value(col, label, curr, prev, value_formatter, delta_formatter=None, inverse=False):
            def _fmt(v):
                return value_formatter(v) if v is not None and not (isinstance(v, float) and np.isnan(v)) else "—"
            if prev is None or (isinstance(prev, float) and np.isnan(prev)):
                col.metric(label=label, value=_fmt(curr), delta="")
                monthly_cards.append({"label": label, "value": _fmt(curr), "delta": "", "delta_color": "normal"})
                return
            if curr is None or (isinstance(curr, float) and np.isnan(curr)):
                col.metric(label=label, value="—", delta="")
                monthly_cards.append({"label": label, "value": "—", "delta": "", "delta_color": "normal"})
                return
            delta = curr - prev
            delta_txt = delta_formatter(delta) if delta_formatter else f"{delta:+,.0f}"
            delta_color = "inverse" if inverse else "normal"
            col.metric(label=label, value=_fmt(curr), delta=delta_txt, delta_color=delta_color)
            monthly_cards.append({"label": label, "value": _fmt(curr), "delta": delta_txt, "delta_color": delta_color})

        metric_delta_value(cA, "Total Generation (kWh)", kwh_curr, kwh_prev, f0, fd0, inverse=False)
        metric_delta_value(cB, "Total Energy Expense (₨)", exp_curr_val, exp_prev_val, f0, fd0, inverse=True)
        metric_delta_value(cC, "Average Cost (₨/kWh)", cpk_curr_val, cpk_prev_val, f2, fd2, inverse=True)
        metric_delta_value(cD, "Solar Share (%)", mix_curr, mix_prev, f1p, fd1p, inverse=False)
        
        # Capture monthly cards for PDF
        capture_metric_cards(monthly_cards, "Monthly Highlights")

    # Composition
    if selector_months:
        section_title("Total Generation by Source", level=2)
        if 'sel_year' not in locals() or sel_year is None:
            sel_year = selector_months[-1].year
        y_mask_year = (months.dt.year == sel_year)

        def _sum_year(s: Optional[pd.Series]) -> float:
            if s is None: return 0.0
            vals = s.iloc[y_mask_year.values]
            vals = vals[vals.notna()]
            return float(vals.sum(skipna=True)) if len(vals) else 0.0

        parts = {
            "Rental": _sum_year(rental_al),
            "PETPAK #01": _sum_year(petpak1_al),
            "GPAK #01": _sum_year(gpak1_al),
            "GPAK #02": _sum_year(gpak2_al),
            "Solar": _sum_year(solar_aligned),
            "LESCO": _sum_year(lesco_aligned),
        }
        labels_vals = [(k, v) for k, v in parts.items() if v and not np.isnan(v) and v > 0]
        if labels_vals:
            labels_all, values_all = zip(*labels_vals)
            df_d_full = pd.DataFrame({"Source": labels_all, "kWh": values_all})
            sel_sources = st.multiselect("Sources to include", options=list(labels_all), default=list(labels_all), key="overview_sources_pick")
            df_d = df_d_full[df_d_full["Source"].isin(sel_sources)]

            comp_opts = ["Donut", "Bar"]
            ctype = st.selectbox("Composition chart type", comp_opts, index=0, key="overview_comp_layout")

            default_color_map = {
                "Rental": RENTAL_GREEN,
                "PETPAK #01": ORANGE_MAIN,
                "GPAK #01": BLUE_DARK,
                "GPAK #02": BLUE_LIGHT,
                "Solar": SOLAR_GOLD,
                "LESCO": LESCO_TEAL,
            }
            base_palette = [default_color_map.get(s, PURPLE) for s in sel_sources]
            palette, custom_map = color_controls(sel_sources, "overview_comp_colors", base_palette)
            color_map_use = {s: (custom_map.get(s) or default_color_map.get(s, PURPLE)) for s in sel_sources}

            if ctype == "Donut":
                fig_d = px.pie(df_d, names="Source", values="kWh", hole=0.55,
                               color="Source", color_discrete_map=color_map_use,
                               title=f"Total Generation — {sel_year}")
                ordered_colors = [color_map_use[s] for s in df_d["Source"]]
                ins_colors = ["#FFFFFF" if (int(c[1:3],16)*0.2126 + int(c[3:5],16)*0.7152 + int(c[5:7],16)*0.0722)/255 < 0.5 else "#111827" for c in ordered_colors]
                fig_d.update_traces(
                    textposition="inside", 
                    textinfo="label+percent",
                    insidetextfont=dict(color=ins_colors),
                    outsidetextfont=dict(color="#111827"),
                    hovertemplate="<b>%{label}</b><br>Generation: %{value:,.0f} kWh<br>Share: %{percent}<br>Year: {sel_year}<extra></extra>"
                )
            else:
                fig_d = px.bar(df_d, x="Source", y="kWh", color="Source",
                               color_discrete_map=color_map_use, title=f"Total Generation — {sel_year}")
                _add_bar_value_labels(fig_d)
            _apply_common_layout(fig_d, f"Total Generation — {sel_year}")
            render_fig(fig_d, key="overview_generation")

    # ENERGY MIX
    if selector_months:
        section_title("Energy Mix by Month", level=2)
        if 'sel_year' not in locals() or sel_year is None:
            sel_year = selector_months[-1].year

        df_mix = pd.DataFrame(index=months)
        if gas_total is not None:   df_mix["Gas"]   = _align(gas_total)
        if lesco_units is not None: df_mix["LESCO"] = _align(lesco_units)
        if solar_total is not None: df_mix["Solar"] = _align(solar_total)
        df_mix = df_mix.loc[[m for m in df_mix.index if m.year == sel_year]]
        df_mix = df_mix.replace(0, np.nan).dropna(how="all").fillna(0.0)

        if not df_mix.empty:
            df_plot = df_mix.reset_index().rename(columns={df_mix.reset_index().columns[0]: "Month"})
            ycols_all = [c for c in ["Gas", "LESCO", "Solar"] if c in df_plot.columns]
            ycols = st.multiselect("Energy sources", options=ycols_all, default=ycols_all, key="mix_sources_pick")
            COLOR_MAP_ENERGY = {"Gas": "#06B6D4", "LESCO": "#1F2937", "Solar": "#F59E0B"}
            fallback = brand_palette("mix", len(ycols))
            base_pal = [COLOR_MAP_ENERGY.get(n, fallback[i]) for i, n in enumerate(ycols)]
            pal_fixed, color_map_fixed = color_controls(ycols, "mix_colors", base_pal)
            color_map_fixed = {n: color_map_fixed.get(n, base_pal[i]) for i, n in enumerate(ycols)}

            ctype = st.selectbox("Mix chart type",
                                 ["Stacked Bar", "Bar", "Line", "Area (Mountain)"],
                                 index=0, key="mix_layout_overview")
            title_mix = f"Energy Mix — {sel_year}"
            if ctype == "Bar":
                fig_mix = _bar_like_from_wide(df_plot, "Month", ycols, title_mix,
                                              stacked=False, palette=pal_fixed, color_map=color_map_fixed)
            elif ctype == "Stacked Bar":
                fig_mix = _bar_like_from_wide(df_plot, "Month", ycols, title_mix,
                                              stacked=True, palette=pal_fixed, color_map=color_map_fixed)
                totals = df_plot[ycols].sum(axis=1).values.astype(float)
                _overlay_totals_text(fig_mix, x=df_plot["Month"], totals=totals)
                _add_visibility_dropdown_for_totals(fig_mix, df_plot, ycols, len(fig_mix.data)-1)
            elif ctype == "Line":
                fig_mix = _line_from_wide(df_plot, "Month", ycols, title_mix, palette=pal_fixed)
            else:
                fig_mix = _area_from_wide(df_plot, "Month", ycols, title_mix, palette=pal_fixed)

            fig_mix.update_yaxes(title_text="kWh", ticksuffix=" kWh", showgrid=True, gridcolor="#e5e7eb", nticks=12)
            fig_mix.update_layout(height=520)
            render_fig(fig_mix, key="overview_energy_mix")

    st.divider()
    # Download options moved to Report tab for better performance

# ───────────────────────────────────────────────────────────────────────────────
# ENERGY SOURCES
# ───────────────────────────────────────────────────────────────────────────────
with tab_sources:
    begin_section("Energy Sources")
    section_title("Engines & Sources", level=2)

    def _engine_block(title, df_eng: Optional[pd.DataFrame], key_prefix: str, context: str):
        if df_eng is None or df_eng.empty:
            return
        df_e = df_eng.copy()

        # UNITS (KWH) — single series with color control
        if "UNITS (KWH)" in df_e.columns:
            df_b = df_e[["Month", "UNITS (KWH)"]].copy().replace(0, np.nan).dropna()
            if not df_b.empty:
                layout = st.selectbox(f"{title}: layout", ["Bar","Line","Area (Mountain)","Lollipop"], index=0, key=f"{key_prefix}_kwh_layout")
                base_pal = brand_palette(context, 1)
                pal, _map = color_controls(["kWh"], f"{key_prefix}_kwh_colors", base_pal)
                if layout == "Bar":
                    fig = _bar_like_from_wide(df_b, "Month", "UNITS (KWH)", f"{title} — kWh", stacked=False, palette=pal)
                elif layout == "Line":
                    fig = _line_from_wide(df_b.rename(columns={"UNITS (KWH)": f"{title} — kWh"}), "Month", [f"{title} — kWh"], f"{title} — kWh", palette=pal)
                elif layout == "Area (Mountain)":
                    fig = _area_from_wide(df_b.rename(columns={"UNITS (KWH)": f"{title} — kWh"}), "Month", [f"{title} — kWh"], f"{title} — kWh", palette=pal)
                else:
                    fig = _lollipop_from_series(df_b.rename(columns={"UNITS (KWH)":"kWh"}), "Month", "kWh", f"{title} — kWh", palette=pal)
                fig.update_yaxes(title_text="kWh", ticksuffix=" kWh")
                _apply_common_layout(fig, f"{title} — kWh")
                render_fig(fig, key=f"sources_{title.lower().replace(' ', '_')}")

        # Efficiency (M3/KWH or kWh/MMBtu) — series multiselect + colors
        y2_candidates = [c for c in ["M3/KWH", "kWh/MMBtu"] if c in df_e.columns]
        if y2_candidates:
            df_l = df_e[["Month"] + y2_candidates].copy()
            for c in y2_candidates: df_l[c] = pd.to_numeric(df_l[c], errors="coerce")
            df_l = df_l.replace(0, np.nan).dropna(subset=y2_candidates, how="all")
            if not df_l.empty:
                picked = st.multiselect(f"{title}: efficiency metrics", options=y2_candidates, default=y2_candidates, key=f"{key_prefix}_eff_pick")
                if picked:
                    layout = st.selectbox(f"{title}: efficiency layout", ["Line","Area (Mountain)","Bar","Stacked Bar","Lollipop"], index=0, key=f"{key_prefix}_eff_layout")
                    base_pal = brand_palette(context, len(picked))
                    pal, _map = color_controls(picked, f"{key_prefix}_eff_colors", base_pal)
                    if layout == "Line":
                        fig2 = _line_from_wide(df_l, "Month", picked, title=f"{title} — Efficiency", palette=pal)
                    elif layout == "Area (Mountain)":
                        fig2 = _area_from_wide(df_l, "Month", picked, title=f"{title} — Efficiency", palette=pal)
                    elif layout == "Bar":
                        fig2 = _bar_like_from_wide(df_l, "Month", picked, f"{title} — Efficiency", stacked=False, palette=pal)
                    elif layout == "Stacked Bar" and len(picked) > 1:
                        fig2 = _bar_like_from_wide(df_l, "Month", picked, f"{title} — Efficiency", stacked=True, palette=pal)
                    else:
                        metric0 = picked[0]
                        fig2 = _lollipop_from_series(df_l.rename(columns={metric0:"Value"}), "Month", "Value", f"{title} — {metric0}", palette=pal[:1])
                    render_fig(fig2, key=f"sources_{title.lower().replace(' ', '_')}_efficiency")

    _engine_block("PETPAK Engine 01", blocks.get("PETPAK ENGINE #01"), "petpak1", context="petpak_engine")
    _engine_block("GPAK Engine 01", blocks.get("GPAK ENGINE #01"), "gpak1", context="gpak_engine_1")
    _engine_block("GPAK Engine 02", blocks.get("GPAK ENGINE #2"),  "gpak2", context="gpak_engine_2")

    # Rental Engine — single series color
    rental_block = blocks.get("RENTAL ENGINE")
    if rental_block is not None and "UNITS (KWH)" in rental_block.columns:
        df_r = rental_block.copy()[["Month", "UNITS (KWH)"]].replace(0, np.nan).dropna()
        if not df_r.empty:
            layout = st.selectbox("Rental Engine: layout", ["Bar","Line","Area (Mountain)","Lollipop"], index=0, key="rental_kwh_layout")
            pal, _ = color_controls(["kWh"], "rental_kwh_colors", [RENTAL_GREEN])
            if layout == "Bar":
                fig = _bar_like_from_wide(df_r, "Month", "UNITS (KWH)", "Rental Engine — kWh", stacked=False, palette=pal)
            elif layout == "Line":
                fig = _line_from_wide(df_r.rename(columns={"UNITS (KWH)":"kWh"}), "Month", ["kWh"], "Rental Engine — kWh", palette=pal)
            elif layout == "Area (Mountain)":
                fig = _area_from_wide(df_r.rename(columns={"UNITS (KWH)":"kWh"}), "Month", ["kWh"], "Rental Engine — kWh", palette=pal)
            else:
                fig = _lollipop_from_series(df_r.rename(columns={"UNITS (KWH)":"kWh"}), "Month", "kWh", "Rental Engine — kWh", palette=pal)
            fig.update_yaxes(title_text="kWh", ticksuffix=" kWh")
            _apply_common_layout(fig, "Rental Engine — kWh")
            render_fig(fig, key="sources_rental_engine")

    # Solar: series multiselect + colors
    solar_block = blocks.get("SOLAR GENERATION")
    if solar_block is not None and not solar_block.empty:
        df_s = solar_block.copy()
        ycols_all = [c for c in ["PETPAK", "GPAK", "TOTAL"] if c in df_s.columns]
        if ycols_all:
            pick = st.multiselect("Solar series", options=ycols_all, default=[c for c in ycols_all if c != "TOTAL"] or ["TOTAL"], key="solar_series_pick")
            df_s2 = df_s[["Month"] + pick].replace(0, np.nan).dropna(subset=pick, how="all").fillna(0.0)
            if not df_s2.empty:
                layout = st.selectbox("Solar layout", ["Stacked Bar","Bar","Line","Area (Mountain)"], index=0, key="solar_layout")
                base_pal = brand_palette("solar_pg" if ("PETPAK" in pick or "GPAK" in pick) else "mix", len(pick))
                pal, _ = color_controls(pick, "solar_colors", base_pal)
                if layout == "Bar":
                    fig = _bar_like_from_wide(df_s2, "Month", pick, "Solar Generation (kWh)", stacked=False, palette=pal)
                elif layout == "Stacked Bar" and len(pick) > 1:
                    fig = _bar_like_from_wide(df_s2, "Month", pick, "Solar Generation (kWh)", stacked=True, palette=pal)
                elif layout == "Line":
                    fig = _line_from_wide(df_s2, "Month", pick, "Solar Generation (kWh)", palette=pal)
                else:
                    fig = _area_from_wide(df_s2, "Month", pick, "Solar Generation (kWh)", palette=pal)
                totals_here = df_s2[pick].sum(axis=1).values
                if layout in ("Bar", "Stacked Bar"):
                    _overlay_totals_text(fig, x=df_s2["Month"], totals=totals_here)
                fig.update_yaxes(title_text="kWh", ticksuffix=" kWh")
                _apply_common_layout(fig, "Solar Generation (kWh)")
                render_fig(fig, key="sources_solar_generation")

    # LESCO — single series color
    lesco_block = blocks.get("LESCO GENERATION")
    if lesco_block is not None and "UNITS (KWH)" in lesco_block.columns:
        df_l = lesco_block.copy()[["Month", "UNITS (KWH)"]].replace(0, np.nan).dropna()
        if not df_l.empty:
            layout = st.selectbox("LESCO layout", ["Bar","Line","Area (Mountain)","Lollipop"], index=0, key="lesco_layout")
            pal, _ = color_controls(["kWh"], "lesco_colors", brand_palette("lesco", 1))
            if layout == "Bar":
                fig = _bar_like_from_wide(df_l, "Month", "UNITS (KWH)", "LESCO — kWh", stacked=False, palette=pal)
            elif layout == "Line":
                fig = _line_from_wide(df_l.rename(columns={"UNITS (KWH)":"LESCO kWh"}), "Month", ["LESCO kWh"], "LESCO — kWh", palette=pal)
            elif layout == "Area (Mountain)":
                fig = _area_from_wide(df_l, "Month", ["LESCO kWh"], "LESCO — kWh", palette=pal)
            else:
                fig = _lollipop_from_series(df_l.rename(columns={"UNITS (KWH)":"kWh"}), "Month", "kWh", "LESCO — kWh", palette=pal)
            fig.update_yaxes(title_text="kWh", ticksuffix=" kWh")
            _apply_common_layout(fig, "LESCO — kWh")
            render_fig(fig, key="sources_lesco")

    st.divider()
    # Download options moved to Report tab for better performance

# ───────────────────────────────────────────────────────────────────────────────
# SOLAR SAVINGS  (LESCO excluded — GAS only)  +  SOLAR PRODUCTION (green)
# ───────────────────────────────────────────────────────────────────────────────
with tab_savings:
    begin_section("Solar Savings")  # IMPORTANT: section key
    section_title("Monthly and Cumulative Savings from Solar (Blended LESCO + GAS)", level=2)

    sv = compute_solar_savings_blended(df_full)
    if sv.empty:
        st.info("Solar data not found.")
    else:
        cols_needed = ["Month", "Solar_kWh", "Saving_vs_GAS", "GAS_rate_used"]
        have = [c for c in cols_needed if c in sv.columns]
        sv2 = sv[have].copy()

        if "Solar_kWh" in sv2.columns:
            sv2 = sv2[sv2["Solar_kWh"].replace(0, np.nan).notna()]

        if sv2.empty or "Saving_vs_GAS" not in sv2.columns:
            st.info("Not enough data to compute blended savings (LESCO + GAS).")
        else:
            sel_m = st.selectbox(
                "Select month",
                options=list(sv2["Month"]),
                index=len(sv2) - 1,
                format_func=lambda d: d.strftime("%b %Y"),
                key="sv_month_selector_gas_only",  # keep your key
            )

            row_idx = sv2.index[sv2["Month"] == sel_m][0]
            row = sv2.loc[row_idx]
            prev_idx = max(row_idx - 1, sv2.index.min())
            prev_row = sv2.iloc[prev_idx] if prev_idx >= sv2.index.min() else None

            c1, c2 = st.columns(2)
            savings_cards = []

            def metric_with_delta(col, label, curr, prev):
                if prev is None or np.isnan(prev):
                    col.metric(label, f"{curr:,.0f}")
                    savings_cards.append({"label": label, "value": f"{curr:,.0f}", "delta": "", "delta_color": "normal"})
                else:
                    delta_val = curr - prev
                    col.metric(label, f"{curr:,.0f}", delta=f"{delta_val:+,.0f}")
                    savings_cards.append({"label": label, "value": f"{curr:,.0f}", "delta": f"{delta_val:+,.0f}", "delta_color": "normal"})

            metric_with_delta(
                c1,
                "Solar Savings (₨, blended)",
                float(row["Saving_vs_GAS"]),
                float(prev_row["Saving_vs_GAS"]) if prev_row is not None else np.nan,
            )
            metric_with_delta(
                c2,
                "Solar Production (kWh)",
                float(row.get("Solar_kWh", np.nan)),
                float(prev_row["Solar_kWh"]) if prev_row is not None else np.nan,
            )
            
            # Capture savings cards for PDF
            capture_metric_cards(savings_cards, "Solar Savings & Production")

            # Monthly savings (BLENDED)
            df_m = sv2.rename(columns={"Saving_vs_GAS": "Savings (₨)"})
            layout_m = st.selectbox(
                "Monthly savings chart",
                ["Line", "Area (Mountain)", "Bar"],
                index=0,
                key="sv_month_layout_gas_only",  # keep your key
            )
            ycols = ["Savings (₨)"]
            pal_savings, _ = color_controls(ycols, "sv_month_colors", [SOLAR_GOLD])

            if layout_m == "Bar":
                sfig = _bar_like_from_wide(df_m, "Month", ycols, "Solar Savings — Monthly (Blended LESCO+GAS)", stacked=False, palette=pal_savings)
            elif layout_m == "Area (Mountain)":
                sfig = _area_from_wide(df_m, "Month", ycols, "Solar Savings — Monthly (Blended LESCO+GAS)", palette=pal_savings)
            else:
                sfig = _line_from_wide(df_m, "Month", ycols, "Solar Savings — Monthly (Blended LESCO+GAS)", palette=pal_savings)
            render_fig(sfig, key="savings_monthly_blended")

            # Cumulative savings (BLENDED)
            sv_cum = sv2.copy()
            sv_cum["Saving_vs_GAS"] = sv_cum["Saving_vs_GAS"].cumsum()
            df_c = sv_cum.rename(columns={"Saving_vs_GAS": "Savings (₨)"})
            layout_c = st.selectbox(
                "Cumulative savings chart",
                ["Bar", "Area (Mountain)", "Line"],
                index=0,
                key="sv_cum_layout_gas_only",  # keep your key
            )
            pal_savings_c, _ = color_controls(ycols, "sv_cum_colors", [SOLAR_GOLD])
            if layout_c == "Bar":
                sfig2 = _bar_like_from_wide(df_c, "Month", ycols, "Solar Savings — Cumulative (Blended LESCO+GAS)", stacked=False, palette=pal_savings_c)
            elif layout_c == "Area (Mountain)":
                sfig2 = _area_from_wide(df_c, "Month", ycols, "Solar Savings — Cumulative (Blended LESCO+GAS)", palette=pal_savings_c)
            else:
                sfig2 = _line_from_wide(df_c, "Month", ycols, "Solar Savings — Cumulative (Blended LESCO+GAS)", palette=pal_savings_c)
            render_fig(sfig2, key="savings_cumulative_blended")

            # SOLAR PRODUCTION
            section_title("Solar Production", level=3)
            df_prod_m = sv2[["Month", "Solar_kWh"]].rename(columns={"Solar_kWh": "Solar Production (kWh)"})
            prod_layout_m = st.selectbox(
                "Monthly solar production chart",
                ["Line", "Area (Mountain)", "Bar"],
                index=0,
                key="solar_prod_month_layout",
            )
            prod_y = ["Solar Production (kWh)"]
            pal_prod, _ = color_controls(prod_y, "solar_prod_colors", [EMERALD])

            if prod_layout_m == "Bar":
                pfig = _bar_like_from_wide(df_prod_m, "Month", prod_y, "Solar Production — Monthly", stacked=False, palette=pal_prod)
            elif prod_layout_m == "Area (Mountain)":
                pfig = _area_from_wide(df_prod_m, "Month", prod_y, "Solar Production — Monthly", palette=pal_prod)
            else:
                pfig = _line_from_wide(df_prod_m, "Month", prod_y, "Solar Production — Monthly", palette=pal_prod)
            render_fig(pfig, key="savings_solar_production_monthly")

            sv_prod_cum = sv2[["Month", "Solar_kWh"]].copy()
            sv_prod_cum["Solar_kWh"] = sv_prod_cum["Solar_kWh"].cumsum()
            df_prod_c = sv_prod_cum.rename(columns={"Solar_kWh": "Solar Production (kWh)"})
            prod_layout_c = st.selectbox(
                "Cumulative solar production chart",
                ["Bar", "Area (Mountain)", "Line"],
                index=0,
                key="solar_prod_cum_layout",
            )
            pal_prod_c, _ = color_controls(prod_y, "solar_prod_cum_colors", [EMERALD])
            if prod_layout_c == "Bar":
                pfig2 = _bar_like_from_wide(df_prod_c, "Month", prod_y, "Solar Production — Cumulative", stacked=False, palette=pal_prod_c)
            elif prod_layout_c == "Area (Mountain)":
                pfig2 = _area_from_wide(df_prod_c, "Month", prod_y, "Solar Production — Cumulative", palette=pal_prod_c)
            else:
                pfig2 = _line_from_wide(df_prod_c, "Month", prod_y, "Solar Production — Cumulative", palette=pal_prod_c)
            render_fig(pfig2, key="savings_solar_production_cumulative")

            # EXPLANATION (BLENDED METHOD)
            section_title("How these savings are calculated (Blended LESCO + GAS)", level=3)
            st.markdown(
                """
                **Formula**  
                • Solar savings = Solar_kWh × ((LESCO cost + GAS cost) ÷ (LESCO units + GAS units))
                """
            )
            if "GAS_rate_used" in row and not np.isnan(row["GAS_rate_used"]):
                st.markdown(f"**Blended rate used (LESCO + GAS):**  **₨ {row['GAS_rate_used']:.4f}/kWh**")

            st.markdown(
                f"**Selected month ({sel_m.strftime('%b %Y')}):**  "
                f"• Solar_kWh = **{row.get('Solar_kWh', float('nan')):,.0f}**  "
                f"• Blended rate = **{row.get('GAS_rate_used', float('nan')):.4f}**  "
                f"⇒ Solar savings = **₨ {row.get('Saving_vs_GAS', 0):,.0f}**"
            )

            with st.expander("Show full savings data table (blended)", expanded=False):
                st.dataframe(sv2)

    st.divider()
    # Download options moved to Report tab for better performance

# ───────────────────────────────────────────────────────────────────────────────
# EXPENSES
# ───────────────────────────────────────────────────────────────────────────────
with tab_expense:
    begin_section("Expenses")
    section_title("Expenses & Unit Costs", level=2)
    df_exp = pd.DataFrame({"Month": months})
    if expense_lesco is not None: df_exp["LESCO Bill (₨)"] = expense_lesco.values
    if expense_gas is not None:   df_exp["Gas Bill (₨)"] = expense_gas.values
    if expense_total is not None: df_exp["Total Expense (₨)"] = expense_total.values
    ycols_all = [c for c in df_exp.columns if c != "Month"]
    if ycols_all:
        ycols_pick = st.multiselect(
            "Expense series",
            options=ycols_all,
            default=ycols_all,
            key="expense_series_pick",
        )
        if ycols_pick:
            df_exp2 = df_exp.copy()
            df_exp2[ycols_pick] = df_exp2[ycols_pick].replace(0, np.nan)
            df_exp2 = df_exp2.dropna(subset=ycols_pick, how="all")
            if not df_exp2.empty:
                layout = st.selectbox(
                    "Expense chart layout",
                    ["Stacked Bar","Bar","Line","Area (Mountain)"],
                    index=0,
                    key="expense_layout",
                )
                default_map = {"LESCO Bill (₨)": LESCO_TEAL, "Gas Bill (₨)": GAS_CRIMSON, "Total Expense (₨)": SLATE}
                base_palette = [default_map.get(c, PURPLE) for c in ycols_pick]
                pal, custom_map = color_controls(ycols_pick, "expense_colors", base_palette)
                color_map = {c: custom_map.get(c, default_map.get(c, pal[i % len(pal)])) for i, c in enumerate(ycols_pick)}

                if layout == "Bar":
                    fig_e = _bar_like_from_wide(df_exp2, "Month", ycols_pick, "Monthly Expenses (₨)", stacked=False, palette=pal, color_map=color_map)
                elif layout == "Stacked Bar" and len(ycols_pick) > 1:
                    fig_e = _bar_like_from_wide(df_exp2, "Month", ycols_pick, "Monthly Expenses (₨)", stacked=True, palette=pal, color_map=color_map)
                    _overlay_totals_text(fig_e, x=df_exp2["Month"], totals=df_exp2[ycols_pick].sum(axis=1).values)
                elif layout == "Line":
                    fig_e = _line_from_wide(df_exp2, "Month", ycols_pick, "Monthly Expenses (₨)", palette=pal)
                else:
                    fig_e = _area_from_wide(df_exp2, "Month", ycols_pick, "Monthly Expenses (₨)", palette=pal)
                _apply_common_layout(fig_e, "Monthly Expenses (₨)")
                render_fig(fig_e, key="expenses_monthly")

    if pkr_per_kwh is not None:
        pk_df = pd.DataFrame({"Month": months, "Avg Cost (₨/kWh)": pkr_per_kwh.values}).replace(0, np.nan).dropna()
        if not pk_df.empty:
            layout_c = st.selectbox("Avg cost chart layout", ["Line","Area (Mountain)","Bar","Lollipop"], index=0, key="avgcost_layout")
            pal_cost, _ = color_controls(["Avg Cost (₨/kWh)"], "avgcost_colors", ["#DC2626"])
            if layout_c == "Bar":
                fig_c = _bar_like_from_wide(pk_df, "Month", "Avg Cost (₨/kWh)", "Overall Cost (₨/kWh)", stacked=False, palette=pal_cost)
            elif layout_c == "Area (Mountain)":
                fig_c = _area_from_wide(pk_df, "Month", ["Avg Cost (₨/kWh)"], "Overall Cost (₨/kWh)", palette=pal_cost)
            elif layout_c == "Lollipop":
                fig_c = _lollipop_from_series(pk_df.rename(columns={"Avg Cost (₨/kWh)":"Value"}), "Month", "Value", "Overall Cost (₨/kWh)", palette=pal_cost)
            else:
                fig_c = _line_from_wide(pk_df, "Month", ["Avg Cost (₨/kWh)"], "Overall Cost (₨/kWh)", palette=pal_cost)
            render_fig(fig_c, key="expenses_avg_cost")

    st.divider()
    # Download options moved to Report tab for better performance

# ───────────────────────────────────────────────────────────────────────────────
# PRODUCTION VS CONSUMPTION
# ───────────────────────────────────────────────────────────────────────────────
with tab_prodcons:
    begin_section("Production vs Consumption")
    section_title("Production vs Energy Consumption", level=2)

    def pick_col_name(cols: List[str], candidates: List[str]) -> Optional[str]:
        lower_map = {c.lower(): c for c in cols}
        for cand in candidates:
            if cand in cols:
                return cand
            k = cand.lower()
            if k in lower_map:
                return lower_map[k]
        return None

    def render_prod_cons(block_df: Optional[pd.DataFrame], title_prefix: str, brand: str):
        if block_df is None or block_df.empty:
            st.info(f"{title_prefix}: data not found.")
            return

        dfp = block_df.copy()
        cols = list(dfp.columns)
        c_prod  = pick_col_name(cols, ["PRODUCTION (KG)", "Production (KG)", "Production KG", "PRODUCTION"])
        c_kwh   = pick_col_name(cols, ["UNITS CONSUMPTION (KWh)", "UNITS CONSUMPTION (KWH)", "Energy Consumption (KWh)", "KWh Consumption", "UNITS (KWH)"])
        c_eff   = pick_col_name(cols, ["KWh/KG", "kWh/Kg", "KWH/KG"])
        c_pkrkg = pick_col_name(cols, ["PKR/KG", "Rs/KG", "Rs per KG", "Cost per KG", "₨/KG"])

        if c_kwh and c_prod:
            df1 = dfp[["Month", c_kwh, c_prod]].copy().replace(0, np.nan).dropna(how="all", subset=[c_kwh, c_prod])
            if not df1.empty:
                layout = st.selectbox(f"{title_prefix}: Prod vs Energy layout", ["Line","Bar","Area (Mountain)"], index=0, key=f"{title_prefix}_pvslayout")
                labels = ["Energy (kWh)", "Production (kg)"]
                base_bar = ORANGE_MAIN if brand == "PETPAK" else BLUE_MAIN
                base_line = RENTAL_GREEN
                pal, cmap = color_controls(labels, f"{title_prefix}_pvs_colors", [base_bar, base_line])

                if layout == "Bar":
                    fig_pc = _bar_like_from_wide(
                        df1.rename(columns={c_kwh:"Energy (kWh)", c_prod:"Production (kg)"}),
                        "Month", labels, f"{title_prefix} — Production vs Energy", stacked=False,
                        palette=pal, color_map={labels[0]: pal[0], labels[1]: pal[1]},
                    )
                elif layout == "Area (Mountain)":
                    fig_pc = _area_from_wide(
                        df1.rename(columns={c_kwh:"Energy (kWh)", c_prod:"Production (kg)"}),
                        "Month", labels, f"{title_prefix} — Production vs Energy", palette=pal,
                    )
                else:
                    fig_pc = go.Figure()
                    fig_pc.add_bar(x=df1["Month"], y=df1[c_kwh], name="Energy (kWh)", marker_color=pal[0])
                    _add_bar_value_labels(fig_pc)
                    fig_pc.add_scatter(
                        x=df1["Month"], y=df1[c_prod], name="Production (kg)",
                        mode="lines+markers+text",
                        line=dict(color=pal[1], width=2.6),
                        marker=dict(size=6, color=pal[1]),
                        text=[_value_text(v) for v in df1[c_prod]], textposition="top center",
                        textfont=dict(color=TEXT_PRIMARY, size=11),
                        yaxis="y2", hovertemplate="%{x|%b %Y}<br>%{y:,.4g}"
                    )
                    fig_pc.update_layout(
                        template="plotly_white",
                        title=f"{title_prefix} — Production vs Energy",
                        yaxis=dict(title="kWh", autorange=True, showgrid=True, gridcolor="#e5e7eb"),
                        yaxis2=dict(title="kg", overlaying="y", side="right", autorange=True, showgrid=False),
                        legend_orientation="h", legend_y=-0.18,
                        hovermode="x unified",
                        margin=dict(t=56, b=10, l=10, r=10),
                        plot_bgcolor="#ffffff", paper_bgcolor="#ffffff",
                        font=dict(color=TEXT_PRIMARY, size=13),
                    )
                render_fig(fig_pc, key=f"prodcons_{title_prefix.lower().replace(' ', '_')}")

        ycols_eff_all = [c for c in [c_eff, c_pkrkg] if c]
        if ycols_eff_all:
            df2 = dfp[["Month"] + ycols_eff_all].copy()
            for c in ycols_eff_all:
                df2[c] = pd.to_numeric(df2[c], errors="coerce")
            df2 = df2.replace(0, np.nan).dropna(subset=ycols_eff_all, how="all")
            if not df2.empty:
                picked = st.multiselect(f"{title_prefix}: efficiency metrics", options=ycols_eff_all, default=ycols_eff_all, key=f"{title_prefix}_eff_pick")
                if picked:
                    layout = st.selectbox(f"{title_prefix}: efficiency layout", ["Line","Area (Mountain)","Bar","Stacked Bar","Lollipop"], index=0, key=f"{title_prefix}_efflayout")
                    pal, _ = color_controls(picked, f"{title_prefix}_eff_colors", brand_palette("petpak_engine" if brand == "PETPAK" else "gpak_engine_1", len(picked)))
                    if layout == "Line":
                        fig_eff = _line_from_wide(df2, "Month", picked, f"{title_prefix} — Efficiency", palette=pal)
                    elif layout == "Area (Mountain)":
                        fig_eff = _area_from_wide(df2, "Month", picked, f"{title_prefix} — Efficiency", palette=pal)
                    elif layout == "Bar":
                        fig_eff = _bar_like_from_wide(df2, "Month", picked, f"{title_prefix} — Efficiency", stacked=False, palette=pal)
                    elif layout == "Stacked Bar" and len(picked) > 1:
                        fig_eff = _bar_like_from_wide(df2, "Month", picked, f"{title_prefix} — Efficiency", stacked=True, palette=pal)
                    else:
                        metric0 = picked[0]
                        fig_eff = _lollipop_from_series(df2.rename(columns={metric0: "Value"}), "Month", "Value", f"{title_prefix} — {metric0}", palette=pal[:1])
                    render_fig(fig_eff, key=f"prodcons_{title_prefix.lower().replace(' ', '_')}_efficiency")

    render_prod_cons(blocks.get("PETPAK PVC"), "PETPAK", brand="PETPAK")
    render_prod_cons(blocks.get("GPAK PVC"), "GPAK", brand="GPAK")

    st.divider()
    # Download options moved to Report tab for better performance

# ───────────────────────────────────────────────────────────────────────────────
# GAS CONSUMPTION  (renders ONLY inside this tab)
# ───────────────────────────────────────────────────────────────────────────────
with tab_gas:
    begin_section("Gas Consumption")
    section_title("Gas Consumption & Rate", level=2)

    gas_block = blocks.get("GAS CONSUMPTION")
    if gas_block is None or gas_block.empty:
        st.info("Gas consumption data not found.")
    else:
        df_g = gas_block.copy()
        if "Month" not in df_g.columns:
            df_g.insert(0, "Month", df_full[("Month", "Month")])

        def pick(cols: List[str], candidates: List[str]) -> Optional[str]:
            low = {c.lower(): c for c in cols}
            for cand in candidates:
                if cand in cols:
                    return cand
                lc = cand.lower()
                if lc in low:
                    return low[lc]
            for c in cols:
                cl = str(c).lower()
                for cand in candidates:
                    if cand.lower() in cl:
                        return c
            return None

        known_total_candidates = [
            "TOTAL GAS", "TOTAL GAS (M3)", "TOTAL GAS (SM3)", "TOTAL (MMBTU)", "TOTAL (MMBtu)",
            "TOTAL CONSUMPTION", "TOTAL", "GAS VOLUME", "TOTAL (M3)", "TOTAL (SM3)"
        ]
        known_rate_candidates = [
            "TOTAL GAS RATE", "GAS RATE", "PKR/SM3", "PKR/M3", "PKR/MMBtu", "PKR/MMBTU",
            "Rs per M3", "Rs/M3", "Rs/MMBtu", "RATE"
        ]
        c_total_gas = pick(list(df_g.columns), known_total_candidates)
        c_gas_rate  = pick(list(df_g.columns), known_rate_candidates)

        maybe_use_cols = []
        for col in df_g.columns:
            if col == "Month":
                continue
            cl = str(col).lower()
            if c_total_gas and cl == str(c_total_gas).lower():
                continue
            if c_gas_rate and cl == str(c_gas_rate).lower():
                continue
            if any(token in cl for token in ["mmbtu", "use", "heater", "plant", "generator", "boiler", "power", "treatment", "burner", "gas"]):
                maybe_use_cols.append(col)

        uses_numeric = []
        if maybe_use_cols:
            dfu = df_g[["Month"] + maybe_use_cols].copy()
            for c in maybe_use_cols:
                dfu[c] = pd.to_numeric(dfu[c], errors="coerce")
            dfu = dfu.replace(0, np.nan).dropna(how="all", subset=maybe_use_cols)
            uses_numeric = [c for c in maybe_use_cols if dfu[c].notna().any()]

            if uses_numeric:
                uses_pick = st.multiselect("Gas 'use' series (MMBtu)", options=uses_numeric, default=uses_numeric, key="gas_use_pick")
                if uses_pick:
                    def key(s: str) -> str: return str(s).strip().lower()
                    COLOR_MAP_FIXED = {
                        "power plant (mmbtu)": CYAN,
                        "oil heater gpak (mmbtu)": TEAL,
                        "oil heaters petpak (mmbtu)": ORANGE_MAIN,
                        "steam generator petpak (mmbtu)": SOLAR_GOLD,
                        "flame treatment gpak (mmbtu)": BLUE_LIGHT,
                    }
                    base_pal = [COLOR_MAP_FIXED.get(key(c), brand_palette("mix", len(uses_pick))[i]) for i, c in enumerate(uses_pick)]
                    pal_use, _ = color_controls(uses_pick, "gas_use_colors", base_pal)

                    layout_use = st.selectbox(
                        "Gas by use — chart type",
                        ["Stacked Bar", "Bar", "Line", "Area (Mountain)"],
                        index=0,
                        key="gas_use_layout",
                    )

                    title_use = "Gas Consumption by Use (MMBtu)"
                    df_plot_use = dfu[["Month"] + uses_pick].copy()

                    if layout_use == "Bar":
                        fig_use = _bar_like_from_wide(df_plot_use, "Month", uses_pick, title_use, stacked=False, palette=pal_use)
                        _add_bar_value_labels(fig_use, inside=False)
                    elif layout_use == "Stacked Bar" and len(uses_pick) > 1:
                        fig_use = _bar_like_from_wide(df_plot_use, "Month", uses_pick, title_use, stacked=True, palette=pal_use)
                        _add_bar_value_labels(fig_use, inside=True)
                        totals = df_plot_use[uses_pick].sum(axis=1).values.astype(float)
                        _overlay_totals_text(fig_use, x=df_plot_use["Month"], totals=totals)
                    elif layout_use == "Line":
                        fig_use = _line_from_wide(df_plot_use, "Month", uses_pick, title_use, palette=pal_use)
                    else:
                        fig_use = _area_from_wide(df_plot_use, "Month", uses_pick, title_use, palette=pal_use)

                    fig_use.update_yaxes(title_text="MMBtu", ticksuffix=" MMBtu")
                    _apply_common_layout(fig_use, title_use)
                    render_fig(fig_use, key="gas_usage")

        dfg = None
        if c_total_gas is None and c_gas_rate is None:
            st.warning("No recognizable 'Total Gas' or 'Gas Rate' columns found in the GAS CONSUMPTION block.")
        else:
            keep_cols = ["Month"] + [c for c in [c_total_gas, c_gas_rate] if c]
            dfg = df_g[keep_cols].copy()
            for c in [c_total_gas, c_gas_rate]:
                if c:
                    dfg[c] = pd.to_numeric(dfg[c], errors="coerce")
            dfg = dfg.replace(0, np.nan).dropna(how="all", subset=[c for c in [c_total_gas, c_gas_rate] if c])

            if not dfg.empty:
                labels_tr = []
                if c_total_gas: labels_tr.append(c_total_gas)
                if c_gas_rate:  labels_tr.append(c_gas_rate)
                base = [EMERALD, VIOLET][:len(labels_tr)]
                pal_gr, _ = color_controls(labels_tr, "gas_total_rate_colors", base)

                fig_g = go.Figure()
                if c_total_gas:
                    fig_g.add_bar(
                        x=dfg["Month"], y=dfg[c_total_gas],
                        name=c_total_gas,
                        marker_color=pal_gr[0]
                    )
                    _add_bar_value_labels(fig_g)

                if c_gas_rate:
                    fig_g.add_scatter(
                        x=dfg["Month"], y=dfg[c_gas_rate],
                        name=c_gas_rate,
                        mode="lines+markers+text",
                        line=dict(color=pal_gr[min(1, len(pal_gr)-1)], width=2.6),
                        marker=dict(size=6, color=pal_gr[min(1, len(pal_gr)-1)]),
                        text=[_value_text(v) for v in dfg[c_gas_rate]],
                        textposition="top center",
                        textfont=dict(color=TEXT_PRIMARY, size=11),
                        yaxis="y2",
                        hovertemplate="%{x|%b %Y}<br>%{y:,.4g}"
                    )

                fig_g.update_layout(
                    template="plotly_white",
                    title="Total Gas & Rate",
                    yaxis=dict(title="MMBtu", autorange=True, showgrid=True, gridcolor="#e5e7eb"),
                    yaxis2=dict(title="₨/MMBtu", overlaying="y", side="right", autorange=True, showgrid=False),
                    legend_orientation="h", legend_y=-0.18,
                    hovermode="x unified",
                    margin=dict(t=56, b=10, l=10, r=10),
                    plot_bgcolor="#ffffff", paper_bgcolor="#ffffff",
                    font=dict(color=TEXT_PRIMARY, size=13),
                )
                render_fig(fig_g, key="gas_rate")

        with st.expander("Show gas consumption data table", expanded=False):
            if dfg is not None and not dfg.empty:
                alias = {}
                if c_total_gas: alias[c_total_gas] = "Total MMBtu"
                if c_gas_rate:  alias[c_gas_rate]  = "PKR/MMBtu"
                st.dataframe(dfg.rename(columns=alias))
            else:
                st.info("No gas consumption data available for table.")

    st.divider()
    # Download options moved to Report tab for better performance

# ───────────────────────────────────────────────────────────────────────────────
# FORECASTING (NEW) — Gaussian Naive Bayes-based bucket forecast
# ───────────────────────────────────────────────────────────────────────────────
from sklearn.naive_bayes import GaussianNB

with tab_forecast:
    begin_section("Forecasting")
    section_title("Monthly Forecast — Next Month Only", level=2)

    if not selector_months:
        st.info("No dated rows found. Please ensure the 'Month' column is populated.")
    else:
        # ---------- helpers ----------
        def _safe_align_series(s):
            try:
                return _align(s)
            except Exception:
                return s.reindex(months) if hasattr(s, "reindex") else None

        def _mk_hist_df():
            hist = pd.DataFrame(index=months)
            if gen_aligned is not None:   hist["Total_kWh"] = gen_aligned
            if solar_aligned is not None: hist["Solar_kWh"] = solar_aligned
            if lesco_aligned is not None: hist["LESCO_kWh"] = lesco_aligned
            if gas_total is not None:     hist["Gas_kWh"]   = _safe_align_series(gas_total)
            if rental_al is not None:     hist["Rental_kWh"]= rental_al
            hist = hist.dropna(how="all").fillna(0.0)
            if "Total_kWh" in hist.columns:
                tot = hist["Total_kWh"].replace(0, np.nan)
                for col in ["Solar_kWh","LESCO_kWh","Gas_kWh","Rental_kWh"]:
                    if col in hist.columns:
                        hist[col.replace("_kWh","_Share")] = (hist[col] / tot).clip(0,1)
                hist = hist.dropna(subset=["Total_kWh"])
            return hist

        def _nb_expected_value_from_bins(X, y_cont, x_new):
            """Bin continuous target into quartiles → GaussianNB → expected value via bin means."""
            if len(y_cont) < 4:
                return float(np.mean(y_cont))
            # quartile bins
            qs = np.quantile(y_cont, [0, .25, .5, .75, 1.0])
            # ensure unique
            qs = np.unique(qs)
            if len(qs) < 3:
                return float(np.mean(y_cont))
            y_bin = np.digitize(y_cont, qs[1:-1], right=True)
            clf = GaussianNB()
            clf.fit(X, y_bin)
            proba = clf.predict_proba(np.asarray([x_new]))[0]
            # means per bin
            means = []
            for b in range(proba.shape[0]):
                vals = y_cont[y_bin == b]
                means.append(vals.mean() if len(vals) else np.nan)
            means = np.nan_to_num(means, nan=np.nanmean(y_cont))
            return float(np.dot(proba, means))

        def _metrics_row(total, solar, gas, lesco, rental, label):
            c1,c2,c3,c4,c5 = st.columns(5)
            forecast_cards = []
            
            c1.metric(f"{label}: Total (kWh)", f"{total:,.0f}")
            forecast_cards.append({"label": f"{label}: Total (kWh)", "value": f"{total:,.0f}", "delta": "", "delta_color": "normal"})
            
            c2.metric("Solar (kWh)",  f"{solar:,.0f}", delta=f"{(solar/max(total,1))*100:,.1f}% share")
            forecast_cards.append({"label": "Solar (kWh)", "value": f"{solar:,.0f}", "delta": f"{(solar/max(total,1))*100:,.1f}% share", "delta_color": "normal"})
            
            c3.metric("Gas (kWh)",    f"{gas:,.0f}",   delta=f"{(gas/max(total,1))*100:,.1f}% share")
            forecast_cards.append({"label": "Gas (kWh)", "value": f"{gas:,.0f}", "delta": f"{(gas/max(total,1))*100:,.1f}% share", "delta_color": "normal"})
            
            c4.metric("LESCO (kWh)",  f"{lesco:,.0f}", delta=f"{(lesco/max(total,1))*100:,.1f}% share")
            forecast_cards.append({"label": "LESCO (kWh)", "value": f"{lesco:,.0f}", "delta": f"{(lesco/max(total,1))*100:,.1f}% share", "delta_color": "normal"})
            
            c5.metric("Rental (kWh)", f"{rental:,.0f}",delta=f"{(rental/max(total,1))*100:,.1f}% share")
            forecast_cards.append({"label": "Rental (kWh)", "value": f"{rental:,.0f}", "delta": f"{(rental/max(total,1))*100:,.1f}% share", "delta_color": "normal"})
            
            # Capture forecast cards for PDF
            capture_metric_cards(forecast_cards, "Forecast Values")

        def _table_chart(df_fx, title, cmap):
            st.dataframe(
                df_fx.assign(Share_pct=(df_fx["Share"]*100).round(2)),
                use_container_width=True
            )
            layout = st.selectbox(
                "Chart layout", ["Bar","Stacked Bar","Line","Area","Donut"],
                index=0, key="fx_layout"
            )
            # color pickers
            c1,c2,c3,c4 = st.columns(4)
            solar_c = c1.color_picker("Solar color",  SOLAR_GOLD)
            gas_c   = c2.color_picker("Gas color",    EMERALD)
            lesco_c = c3.color_picker("LESCO color",  LESCO_TEAL)
            rent_c  = c4.color_picker("Rental color", RENTAL_GREEN)
            cmap = {"Solar":solar_c,"Gas":gas_c,"LESCO":lesco_c,"Rental":rent_c}

            if layout in ["Bar","Stacked Bar"]:
                fig = px.bar(df_fx, x="Source", y="kWh", color="Source",
                             title=title, color_discrete_map=cmap)
                if layout == "Stacked Bar":
                    fig.update_traces(offsetgroup=None)  # single stack
                _add_bar_value_labels(fig)
            elif layout == "Line":
                fig = px.line(df_fx, x="Source", y="kWh", color="Source",
                              markers=True, title=title, color_discrete_map=cmap)
            elif layout == "Area":
                fig = px.area(df_fx, x="Source", y="kWh", color="Source",
                              title=title, color_discrete_map=cmap, groupnorm="")
            else:  # Donut
                fig = px.pie(df_fx, names="Source", values="kWh", hole=.5,
                             title=title, color="Source", color_discrete_map=cmap)
                fig.update_traces(
                    textposition="inside", 
                    textinfo="percent+label+value",
                    hovertemplate="<b>%{label}</b><br>Generation: %{value:,.0f} kWh<br>Share: %{percent}<br>Forecast Period<extra></extra>"
                )
            _apply_common_layout(fig, title)
            render_fig(fig)

        # ---------- next month target ----------
        last_dt = pd.to_datetime(selector_months[-1])
        next_month = (last_dt + pd.offsets.MonthBegin(1)).strftime("%b %Y")
        st.markdown(f"**Target Month:** {next_month}")

        hist = _mk_hist_df()
        if hist.empty:
            st.warning("Insufficient history to model.")
            st.stop()

        # base features for NB
        H = hist.reset_index().rename(columns={"index":"Month"}).copy()
        H["m_idx"] = range(len(H))
        # Build t -> t+1 target for Total
        H["Target_Total_next"] = H["Total_kWh"].shift(-1)
        Hm = H.dropna(subset=["Target_Total_next"]).copy()

        # UI: Mode + Sensitivity
        mode = st.radio("Mode", ["Naive Bayes (Auto)","Manual (User-defined)"], horizontal=True, key="fc_mode")
        st.caption("Tip: adjust the sensitivity controls to see how changes ripple through outputs.")

        # Sensitivity controls (apply in both modes)
        s1,s2 = st.columns(2)
        total_nudge = s1.slider("Total Generation nudge (%)", -30, 30, 0, 1,
                                help="Applies AFTER the forecast or manual entry.")
        share_nudge = s2.slider("Shares equalize/tilt (%)", -20, 20, 0, 1,
                                help="Negative pulls shares toward equal; positive magnifies larger shares.")

        # ---------- NAIVE BAYES (Auto) ----------
        if mode.startswith("Naive"):
            # Predict Total_kWh_next using NB on (m_idx, Total_kWh, Solar, Gas, LESCO, Rental)
            cols_feats = [c for c in ["m_idx","Total_kWh","Solar_KWh","Solar_kWh","Gas_KWh","Gas_kWh","LESCO_KWh","LESCO_kWh","Rental_kWh"] if c in Hm.columns]
            # normalize common typos
            if "Solar_KWh" in cols_feats and "Solar_kWh" not in cols_feats: Hm["Solar_kWh"]=Hm["Solar_KWh"]
            if "Gas_KWh" in cols_feats and "Gas_kWh" not in cols_feats:     Hm["Gas_kWh"]=Hm["Gas_KWh"]
            if "LESCO_KWh" in cols_feats and "LESCO_kWh" not in cols_feats: Hm["LESCO_kWh"]=Hm["LESCO_KWh"]
            cols_feats = [c for c in ["m_idx","Total_kWh","Solar_kWh","Gas_kWh","LESCO_kWh","Rental_kWh"] if c in Hm.columns]
            X = Hm[cols_feats].values
            y = Hm["Target_Total_next"].values.astype(float)

            # Expected value via probabilistic bins
            x_new = np.array([
                H["m_idx"].max()+1,
                H["Total_kWh"].iloc[-1],
                H.get("Solar_kWh", pd.Series([0])).iloc[-1] if "Solar_kWh" in H else 0.0,
                H.get("Gas_kWh",   pd.Series([0])).iloc[-1] if "Gas_kWh" in H   else 0.0,
                H.get("LESCO_kWh", pd.Series([0])).iloc[-1] if "LESCO_kWh" in H else 0.0,
                H.get("Rental_kWh",pd.Series([0])).iloc[-1] if "Rental_kWh" in H else 0.0,
            ][:len(cols_feats)])

            total_nb = _nb_expected_value_from_bins(X, y, x_new)

            # Predict shares (as before) with GaussianNB on binned shares
            def bin_share(s): return pd.cut(s, bins=[-0.001, 0.333, 0.666, 1.0], labels=[0,1,2]).astype(int)
            share_cols = [c for c in ["Solar_Share","Gas_Share","LESCO_Share","Rental_Share"] if c in H]
            exp_shares = {}
            Xs = H[["m_idx","Total_kWh"]].values
            for col in share_cols:
                yb = bin_share(H[col]).values
                if len(np.unique(yb)) < 2:
                    exp_shares[col] = float(H[col].mean())
                else:
                    clf = GaussianNB()
                    clf.fit(Xs, yb)
                    proba = clf.predict_proba(np.array([[H["m_idx"].max()+1, total_nb]]))[0]
                    mids = np.array([0.165, 0.5, 0.835])
                    exp_shares[col] = float((proba*mids).sum())
            # Normalize
            ssum = sum(exp_shares.values()) or 1.0
            for k in exp_shares: exp_shares[k] /= ssum

            # Apply sensitivity nudges
            total_nb = total_nb * (1 + total_nudge/100.0)
            shares = {
                "Solar":  exp_shares.get("Solar_Share", 0.0),
                "Gas":    exp_shares.get("Gas_Share", 0.0),
                "LESCO":  exp_shares.get("LESCO_Share", 0.0),
                "Rental": exp_shares.get("Rental_Share", 0.0),
            }
            if share_nudge != 0:
                # move toward equal (25%) or magnify disparity
                arr = np.array(list(shares.values()), dtype=float)
                if share_nudge < 0:
                    arr = arr + (0.25 - arr)*abs(share_nudge)/20.0
                else:
                    arr = np.power(arr, 1.0 + share_nudge/20.0)
                    arr = arr / arr.sum()
                shares = dict(zip(shares.keys(), arr))

            solar_kwh  = shares["Solar"]  * total_nb
            gas_kwh    = shares["Gas"]    * total_nb
            lesco_kwh  = shares["LESCO"]  * total_nb
            rental_kwh = shares["Rental"] * total_nb

            section_title("Naive Bayes Forecast — All Values", level=3)
            _metrics_row(total_nb, solar_kwh, gas_kwh, lesco_kwh, rental_kwh, f"{next_month}")

            df_fx = pd.DataFrame([
                ["Solar",  shares["Solar"],  solar_kwh],
                ["Gas",    shares["Gas"],    gas_kwh],
                ["LESCO",  shares["LESCO"],  lesco_kwh],
                ["Rental", shares["Rental"], rental_kwh],
            ], columns=["Source","Share","kWh"])

            _table_chart(df_fx, f"NB Expected Source Consumption — {next_month}", {
                "Solar": SOLAR_GOLD, "Gas": EMERALD, "LESCO": LESCO_TEAL, "Rental": RENTAL_GREEN
            })

            # Optional: Savings/Cost if your app exposes a function. We try to call it; else skip gracefully.
            with st.expander("Costs & Savings (if available)"):
                try:
                    # Example hooks; replace with your app’s real functions if present.
                    if "calc_costs_pkrs" in globals():
                        costs = calc_costs_pkrs(solar_kwh, gas_kwh, lesco_kwh, rental_kwh)
                        st.write(costs)
                    elif "compute_solar_savings" in globals():
                        sv = compute_solar_savings(solar_kwh, baseline="LESCO")
                        st.metric("Estimated Solar Savings (PKR)", f"{sv:,.0f}")
                        # Capture individual metric for PDF
                        capture_metric_cards([{"label": "Estimated Solar Savings (PKR)", "value": f"{sv:,.0f}", "delta": "", "delta_color": "normal"}], "Estimated Solar Savings")
                    else:
                        st.info("Attach your cost/savings function (e.g., calc_costs_pkrs or compute_solar_savings) to display numbers here.")
                except Exception as e:
                    st.info(f"Costs/Savings not computed: {e}")

        # ---------- MANUAL (User-defined) ----------
        else:
            section_title("Manual Forecast — User Inputs", level=3)
            c1,c2 = st.columns(2)
            with c1:
                total_manual = st.number_input("Total Generation (kWh)", min_value=0.0, value=float(H["Total_kWh"].iloc[-1] if len(H) else 0.0), step=1000.0)
                input_mode = st.radio("Input sources as", ["Shares (%)","Absolute kWh"], horizontal=True, key="manual_input_mode")
            with c2:
                st.caption("Enter core values to test scenarios (e.g., tariff changes, availability, outages).")

            if input_mode.startswith("Shares"):
                s1 = st.slider("Solar %",  0, 100, 25, 1)
                s2 = st.slider("Gas %",    0, 100, 25, 1)
                s3 = st.slider("LESCO %",  0, 100, 25, 1)
                s4 = st.slider("Rental %", 0, 100, 25, 1)
                totp = max(1, s1+s2+s3+s4)
                shares = np.array([s1,s2,s3,s4], dtype=float) / totp
                solar_kwh, gas_kwh, lesco_kwh, rental_kwh = shares * total_manual
            else:
                solar_kwh  = st.number_input("Solar (kWh)",  min_value=0.0, value=0.0, step=1000.0)
                gas_kwh    = st.number_input("Gas (kWh)",    min_value=0.0, value=0.0, step=1000.0)
                lesco_kwh  = st.number_input("LESCO (kWh)",  min_value=0.0, value=0.0, step=1000.0)
                rental_kwh = st.number_input("Rental (kWh)", min_value=0.0, value=0.0, step=1000.0)
                total_manual = max(total_manual, solar_kwh + gas_kwh + lesco_kwh + rental_kwh)

            # Apply sensitivity nudges
            total_manual = total_manual * (1 + total_nudge/100.0)
            tot = max(1.0, total_manual)
            shares_manual = {
                "Solar":  solar_kwh/tot,
                "Gas":    gas_kwh/tot,
                "LESCO":  lesco_kwh/tot,
                "Rental": rental_kwh/tot,
            }
            if share_nudge != 0:
                arr = np.array(list(shares_manual.values()), dtype=float)
                if share_nudge < 0:
                    arr = arr + (0.25 - arr)*abs(share_nudge)/20.0
                else:
                    arr = np.power(arr, 1.0 + share_nudge/20.0); arr = arr/arr.sum()
                solar_kwh, gas_kwh, lesco_kwh, rental_kwh = arr * tot

            _metrics_row(tot, solar_kwh, gas_kwh, lesco_kwh, rental_kwh, f"{next_month} (Manual)")

            df_fx = pd.DataFrame([
                ["Solar",  solar_kwh/tot,  solar_kwh],
                ["Gas",    gas_kwh/tot,    gas_kwh],
                ["LESCO",  lesco_kwh/tot,  lesco_kwh],
                ["Rental", rental_kwh/tot, rental_kwh],
            ], columns=["Source","Share","kWh"])

            _table_chart(df_fx, f"Manual Expected Source Consumption — {next_month}", {
                "Solar": SOLAR_GOLD, "Gas": EMERALD, "LESCO": LESCO_TEAL, "Rental": RENTAL_GREEN
            })

            with st.expander("Costs & Savings (if available)"):
                try:
                    if "calc_costs_pkrs" in globals():
                        costs = calc_costs_pkrs(solar_kwh, gas_kwh, lesco_kwh, rental_kwh)
                        st.write(costs)
                    elif "compute_solar_savings" in globals():
                        sv = compute_solar_savings(solar_kwh, baseline="LESCO")
                        st.metric("Estimated Solar Savings (PKR)", f"{sv:,.0f}")
                        # Capture individual metric for PDF
                        capture_metric_cards([{"label": "Estimated Solar Savings (PKR)", "value": f"{sv:,.0f}", "delta": "", "delta_color": "normal"}], "Estimated Solar Savings")
                    else:
                        st.info("Attach your cost/savings function to display numbers here.")
                except Exception as e:
                    st.info(f"Costs/Savings not computed: {e}")

    st.divider()
    # Download options moved to Report tab for better performance

# ───────────────────────────────────────────────────────────────────────────────
# COMPARISON
# ───────────────────────────────────────────────────────────────────────────────
with tab_compare:
    begin_section("Comparison")
    section_title("Custom Comparison", level=2)

    choices_by_unit = {
        "kWh": {
            "Total Generation (kWh)": gen_aligned,
            "LESCO (kWh)": lesco_aligned,
            "Solar (kWh)": solar_aligned,
            "PETPAK Engine #01 (kWh)": petpak1_al,
            "GPAK Engine #01 (kWh)": gpak1_al,
            "GPAK Engine #02 (kWh)": gpak2_al,
            "Rental Engine (kWh)": rental_al,
        },
        "₨": {
            "LESCO Bill (₨)": (pd.Series(expense_lesco.values, index=months) if expense_lesco is not None else None),
            "Gas Bill (₨)":   (pd.Series(expense_gas.values,   index=months) if expense_gas   is not None else None),
            "Total Expense (₨)": (pd.Series(expense_total.values, index=months) if expense_total is not None else None),
        },
        "₨/kWh": {
            "Average Cost (₨/kWh)": (pd.Series(pkr_per_kwh.values, index=months) if pkr_per_kwh is not None else None),
        },
        "MMBtu": {}
    }

    gas_block_for_compare = blocks.get("GAS CONSUMPTION")
    if gas_block_for_compare is not None and not gas_block_for_compare.empty:
        df_gcmp = gas_block_for_compare.copy()
        if "Month" not in df_gcmp.columns:
            df_gcmp.insert(0, "Month", df_full[("Month", "Month")])
        for col in df_gcmp.columns:
            if col == "Month":
                continue
            if "mmbtu" in str(col).lower():
                s = pd.to_numeric(df_gcmp[col], errors="coerce")
                s.index = pd.to_datetime(df_gcmp["Month"])
                s = s.reindex(months.values)
                choices_by_unit["MMBtu"][str(col)] = s

    unit = st.selectbox("Select unit", ["kWh", "MMBtu", "₨", "₨/kWh"], index=0, key="cmp_unit")
    unit_map_all = choices_by_unit.get(unit, {})
    unit_map = {name: s for name, s in unit_map_all.items() if s is not None and s.notna().any()}
    if not unit_map:
        st.info(f"No data available for unit: {unit}")
    else:
        default_labels = list(unit_map.keys())[:2] if len(unit_map) >= 2 else list(unit_map.keys())
        labels = st.multiselect("Select series to compare", options=list(unit_map.keys()), default=default_labels, key="cmp_labels")
        if labels:
            chart_type = st.selectbox(
                "Chart type",
                ["Line", "Bar", "Stacked Bar", "Area (Mountain)", "Donut (by month)", "Donut (year total)"],
                index=0,
                key="cmp_layout"
            )

            base_palette = brand_palette("generic", len(labels))
            palette, color_map = color_controls(labels, "cmp_colors", base_palette)

            df_cmp = pd.DataFrame({"Month": pd.to_datetime(months)})
            for lbl in labels:
                df_cmp[lbl] = unit_map[lbl].values
            df_cmp = df_cmp.replace(0, np.nan).dropna(how="all", subset=labels)

            if not df_cmp.empty:
                title = f"Comparison ({unit})"
                if chart_type in ("Line", "Bar", "Stacked Bar", "Area (Mountain)"):
                    if chart_type == "Line":
                        fig_cmp = _line_from_wide(df_cmp, "Month", labels, title, palette=palette)
                    elif chart_type == "Bar":
                        fig_cmp = _bar_like_from_wide(df_cmp, "Month", labels, title, stacked=False, palette=palette)
                    elif chart_type == "Stacked Bar":
                        fig_cmp = _bar_like_from_wide(df_cmp, "Month", labels, title, stacked=True, palette=palette)
                        totals = df_cmp[labels].sum(axis=1).values.astype(float)
                        _overlay_totals_text(fig_cmp, x=df_cmp["Month"], totals=totals)
                    else:
                        fig_cmp = _area_from_wide(df_cmp, "Month", labels, title, palette=palette)

                    if unit in ("kWh", "MMBtu"):
                        fig_cmp.update_yaxes(title_text=unit, ticksuffix=f" {unit}")

                    _apply_common_layout(fig_cmp, title)
                    render_fig(fig_cmp, key=f"compare_{title.lower().replace(' ', '_').replace('(', '').replace(')', '')}")

                elif chart_type == "Donut (by month)":
                    valid_rows = df_cmp.dropna(subset=labels, how="all")
                    if valid_rows.empty:
                        st.info("No single month has data for the selected series.")
                    else:
                        valid_months = pd.to_datetime(valid_rows["Month"].values)
                        chosen_month = st.selectbox(
                            "Pick month for donut",
                            options=valid_months,
                            index=len(valid_months) - 1,
                            format_func=lambda d: d.strftime("%b %Y"),
                            key="cmp_donut_month"
                        )
                        row = df_cmp.loc[pd.to_datetime(df_cmp["Month"]) == chosen_month, labels].iloc[0]
                        pie_df = pd.DataFrame({"Series": labels, "Value": [row[l] if pd.notna(row[l]) else 0 for l in labels]})
                        pie_df = pie_df[pie_df["Value"] > 0]

                        if not pie_df.empty:
                            color_discrete_map = {lbl: (color_map.get(lbl) or palette[i]) for i, lbl in enumerate(labels)}
                            fig_pie = px.pie(
                                pie_df, names="Series", values="Value", hole=0.55,
                                color="Series", color_discrete_map=color_discrete_map,
                                title=f"{title} — {chosen_month.strftime('%b %Y')}"
                            )
                            ordered_colors = [color_discrete_map[s] for s in pie_df["Series"]]
                            ins_colors = []
                            for c in ordered_colors:
                                try:
                                    r, g, b = int(c[1:3],16), int(c[3:5],16), int(c[5:7],16)
                                    lum = (0.2126*r + 0.7152*g + 0.0722*b)/255
                                    ins_colors.append("#111827" if lum > 0.6 else "#FFFFFF")
                                except Exception:
                                    ins_colors.append("#111827")
                            fig_pie.update_traces(
                                textposition="inside", 
                                textinfo="label+percent",
                                insidetextfont=dict(color=ins_colors),
                                outsidetextfont=dict(color="#111827"),
                                hovertemplate=f"<b>%{{label}}</b><br>Generation: %{{value:,.0f}} {unit}<br>Share: %{{percent}}<br>Month: {chosen_month.strftime('%b %Y')}<extra></extra>"
                            )
                            _apply_common_layout(fig_pie, fig_pie.layout.title.text)
                            render_fig(fig_pie, key=f"compare_{title.lower().replace(' ', '_').replace('(', '').replace(')', '')}_donut")

                with st.expander("Show comparison data table", expanded=False):
                    st.dataframe(df_cmp)

    st.divider()
    # Download options moved to Report tab for better performance

# ───────────────────────────────────────────────────────────────────────────────
# RAW DATA TAB
# ───────────────────────────────────────────────────────────────────────────────
with tab_data:
    begin_section("Data")
    section_title("Raw Data (Sheet: GENERATION AND EXPENSE)", level=2)
    st.dataframe(df_full)

    st.caption(
        "Logo sizing is controlled in-code only. "
        "Edit `LOGO_MAX_HEIGHT_PX`, `LOGO_MAX_WIDTH_PX`, `WM_WIDTH_VW`, and `WM_GAP_VW` at the top of the file."
    )

    st.divider()
    # Download options moved to Report tab for better performance

# ───────────────────────────────────────────────────────────────────────────────
# REPORT — one-page consolidated deck of figures from Overview → Gas Consumption
# ───────────────────────────────────────────────────────────────────────────────
# ───────────────────────────────────────────────────────────────────────────────
# REPORT — one-page consolidated deck of figures from Overview → Gas Consumption
# ───────────────────────────────────────────────────────────────────────────────
with tab_report:
    begin_section("Report")
    section_title("Consolidated Report (Overview → Gas Consumption)", level=2)

    section_order = [
        "Overview",
        "Energy Sources", 
        "Solar Savings",
        "Expenses",
        "Production vs Consumption",
        "Gas Consumption",
        "Comparison",
    ]

    st.caption("This page compiles all charts from the tabs below. Downloads export the whole page in one file.")

    any_content = False
    for sec in section_order:
        figs = SECTION_FIGS.get(sec, [])
        cards_data = SECTION_CARDS.get(sec, [])
        
        st.markdown(f"#### {sec}")
        
        # Show cards if they exist for this section
        if cards_data:
            for card_group in cards_data:
                if not card_group:
                    continue
                
                # Extract cards and title
                if isinstance(card_group, dict) and 'cards' in card_group:
                    cards = card_group['cards']
                    card_section_title = card_group.get('title', '')
                else:
                    cards = card_group
                    card_section_title = ''
                
                if not cards:
                    continue
                
                # Display section title if provided
                if card_section_title:
                    st.markdown(f"**{card_section_title}**")
                
                # Display cards in columns
                cols = st.columns(len(cards))
                for i, card in enumerate(cards):
                    with cols[i]:
                        delta_value = card.get('delta', '')
                        delta_color = card.get('delta_color', 'normal')
                        
                        if delta_value:
                            st.metric(
                                label=card['label'],
                                value=card['value'],
                                delta=delta_value,
                                delta_color=delta_color
                            )
                        else:
                            st.metric(
                                label=card['label'],
                                value=card['value']
                            )
                
                st.write("")  # Add spacing after cards
        
        # Show charts
        if not figs:
            st.info(f"No charts captured yet from **{sec}**. Open that tab once to populate its visuals.")
            st.write("")
            continue

        for i, fig in enumerate(figs, start=1):
            try:
                fig_copy = go.Figure(fig)
            except Exception:
                fig_copy = fig
            if not (getattr(fig_copy, "layout", None) and getattr(fig_copy.layout, "title", None) and getattr(fig_copy.layout.title, "text", None)):
                fig_copy.update_layout(title=f"{sec} — Figure {i}")
            render_fig(fig_copy, key=f"report_{sec.lower().replace(' ', '_')}_{i}")  # captured into SECTION_FIGS["Report"]
            any_content = True

        st.markdown("---")

    if not any_content:
        st.warning("No visuals yet. Open the tabs above first, then return.")

    section_title("Export Consolidated Report", level=2)
    render_export_row("Report", "Consolidated Report — Overview to Gas", "powerhouse_report_all")

# Always show footer at the bottom
render_footer_ui()






